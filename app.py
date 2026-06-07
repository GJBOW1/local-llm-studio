"""Local LLM Studio — a fully offline Flask proxy to a local Ollama server.

Flask sits between the browser and Ollama so the UI never makes a cross-origin
request (no CORS headaches), and matches the user's Flask/Python stack. Everything
runs locally: Flask on :5000, Ollama on :11434. No internet required once the
models are pulled.
"""

from __future__ import annotations

import base64
import json
import os
import re
import secrets
import subprocess
import tempfile
import threading
import time
import uuid
from datetime import datetime, timezone
from typing import Any, Iterator

import requests
from flask import Flask, Response, jsonify, render_template, request, send_file

from mcp_bridge import MCPBridge

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 64 * 1024 * 1024  # 64 MB cap on uploads / image payloads

OLLAMA_HOST: str = os.environ.get("OLLAMA_HOST", "http://localhost:11434").rstrip("/")
EXTRACT_MAX_CHARS: int = int(os.environ.get("LLS_EXTRACT_MAX_CHARS", "20000"))  # doc-text cap (~context-safe)
REQUEST_TIMEOUT: float = 10.0  # seconds, for non-streaming control calls

# Effective context window. Ollama's default num_ctx is small and truncates silently,
# so we set it explicitly (capped for 16 GB RAM) — which also makes the UI's context
# meter truthful, since the displayed window is the one actually in use.
NUM_CTX_CAP: int = int(os.environ.get("LLS_NUM_CTX", "8192"))
_CTX_CACHE: dict[str, int] = {}  # model -> trained context_length (from /api/show)
# How long Ollama keeps the model resident after a request. "30m" keeps it warm
# through a chat session (no ~45s reload on big models like gemma4), then frees the
# RAM when idle. Set "-1" to keep it resident forever, "0" to unload immediately.
KEEP_ALIVE: str = os.environ.get("LLS_KEEP_ALIVE", "30m")

# Saved diagrams/artifacts live as one JSON file each, on disk beside the app —
# so they survive a browser-cache wipe and are real files the user can find.
ARTIFACTS_DIR: str = os.path.join(os.path.dirname(os.path.abspath(__file__)), "artifacts")
# Generated media the user chose to Save is copied here (durable, unlike the temp
# MEDIA_DIR that's wiped on reboot), and referenced by an "embed" artifact.
ARTIFACT_MEDIA_DIR: str = os.path.join(ARTIFACTS_DIR, "media")
ARTIFACT_TYPES: frozenset[str] = frozenset({"mermaid", "svg", "html", "chart", "embed"})
ARTIFACT_ID_RE = re.compile(r"^[a-z0-9][a-z0-9-]{0,63}$")  # blocks path traversal
MAX_ARTIFACT_BYTES: int = 400_000  # generous per-artifact source cap

# Second-brain grounding: shell out to the local skippy-rag index (nomic-embed-text)
# so chats are grounded in the user's vault. Fully local; degrades silently if the
# index or Ollama is unavailable. Everything here is env-overridable.
RAG_DIR: str = os.environ.get("LLS_RAG_DIR", os.path.expanduser("~/skippy-rag"))
RAG_PYTHON: str = os.path.join(RAG_DIR, ".venv", "bin", "python")
RAG_QUERY: str = os.path.join(RAG_DIR, "query.py")
RAG_INDEX: str = os.path.join(RAG_DIR, "index.json")
RAG_GRAPH: str = os.path.join(RAG_DIR, "graphify-out", "graph.json")
RAG_MODEL: str = "nomic-embed-text"  # the actual embed model (truthful health readout)
RAG_TOP_K: int = int(os.environ.get("LLS_RAG_TOP_K", "5"))
RAG_MIN_SCORE: float = float(os.environ.get("LLS_RAG_MIN_SCORE", "0.6"))  # benchmark-tuned (skippy-rag/eval.py): 0.5 fired on every prompt incl. PII notes; 0.6 keeps 100% recall, ~82% specificity
RAG_MAX_CHARS: int = int(os.environ.get("LLS_RAG_MAX_CHARS", "1500"))
RAG_TIMEOUT: float = float(os.environ.get("LLS_RAG_TIMEOUT", "25"))

# The Obsidian vault whose wiki/ Graphify indexed — lets a clicked graph node open
# its source .md in a preview. Override via LLS_VAULT_ROOT.
VAULT_ROOT: str = os.environ.get(
    "LLS_VAULT_ROOT", "~/Obsidian"
)
VAULT_WIKI: str = os.path.join(VAULT_ROOT, "wiki")


def _resolve_note_path(source_file: str) -> str | None:
    """Map a graph node's source_file to a real .md inside the vault wiki, with
    traversal blocked. Graphify stores it prefixed with 'relative/path/' and
    relative to the indexed wiki root. Returns None if not an in-bounds readable note."""
    rel = (source_file or "").strip()
    if rel.startswith("relative/path/"):
        rel = rel[len("relative/path/"):]
    rel = rel.lstrip("/")
    if not rel or not rel.endswith(".md"):
        return None
    root = os.path.realpath(VAULT_WIKI)
    path = os.path.realpath(os.path.join(root, rel))
    if path != root and not path.startswith(root + os.sep):
        return None  # escaped the wiki root
    return path if os.path.isfile(path) else None

# Read-only MCP tool bridge (Monarch, …). Warmed in a background thread so the
# first chat doesn't block on server startup; degrades to no-tools if unavailable.
MCP = MCPBridge()
threading.Thread(target=MCP._ensure_ready, daemon=True).start()
MAX_TOOL_ROUNDS: int = 6


# Local secrets (API keys) from a gitignored config.local.json, env taking priority.
def _load_local_config() -> dict[str, Any]:
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "config.local.json")
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except (OSError, ValueError):
        return {}


_LOCAL_CONFIG = _load_local_config()
BRAVE_KEY: str = os.environ.get("BRAVE_API_KEY") or _LOCAL_CONFIG.get("BRAVE_API_KEY", "")
WEATHER_KEY: str = os.environ.get("OPENWEATHER_API_KEY") or _LOCAL_CONFIG.get("OPENWEATHER_API_KEY", "")
# The RAG retrieval gate is also tunable from config.local.json — a durable, no-code-edit
# knob for a GUI-launched app where shell env vars don't reach the .app. Precedence:
# env var > config.local.json > the code default set above. (Benchmark with skippy-rag/eval.py.)
if "LLS_RAG_MIN_SCORE" not in os.environ and "LLS_RAG_MIN_SCORE" in _LOCAL_CONFIG:
    try:
        RAG_MIN_SCORE = float(_LOCAL_CONFIG["LLS_RAG_MIN_SCORE"])
    except (TypeError, ValueError):
        pass  # keep the code default if the config value is malformed


def _save_local_config_key(name: str, value: str) -> None:
    """Persist (or clear) one key in the gitignored config.local.json, preserving
    the rest. Used by the cloud-provider connect/disconnect endpoints so a key the
    user enters in Settings survives a restart. Never logged."""
    path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "config.local.json")
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, dict):
            data = {}
    except (OSError, ValueError):
        data = {}
    if value:
        data[name] = value
    else:
        data.pop(name, None)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)


# Cloud model providers (opt-in; off by default). Keys live only in env or the
# gitignored config.local.json — never committed, never returned to the client.
PROVIDER_CONFIG_KEY: dict[str, str] = {
    "anthropic": "ANTHROPIC_API_KEY",
    "openai": "OPENAI_API_KEY",
    "gemini": "GEMINI_API_KEY",
    "grok": "XAI_API_KEY",
}
PROVIDER_KEYS: dict[str, str] = {
    p: (os.environ.get(cfg) or _LOCAL_CONFIG.get(cfg, "")) for p, cfg in PROVIDER_CONFIG_KEY.items()
}
PROVIDER_MODELS: dict[str, list[str]] = {p: [] for p in PROVIDER_CONFIG_KEY}

# ── Curated model menu (Phase 6) ──────────────────────────────────────────
# Only these models are offered in the studio, grouped by provider — everything
# else is hidden. EDIT THESE LISTS to change the picks (the menu and the
# auto-download both read from here). Local (Ollama) models are auto-pulled on
# first use; cloud models are filtered to this set of the account's models.
CURATED_LOCAL: dict[str, list[str]] = {
    "Qwen": ["qwen3:8b", "qwen3:14b", "qwen3:32b"],
    "NVIDIA": ["nemotron-mini:4b", "nemotron:70b"],
    "Gemma (Google)": ["gemma4:12b-it-qat", "gemma4:12b", "gemma3:12b-it-qat"],
}
CURATED_CLOUD: dict[str, list[str]] = {
    "anthropic": ["claude-opus-4-8", "claude-sonnet-4-6", "claude-haiku-4-5-20251001"],
    "openai": ["gpt-5.5", "gpt-5.5-pro", "gpt-5.4-mini"],
    "grok": ["grok-4.3", "grok-4.20-0309-reasoning", "grok-4.20-multi-agent-0309"],
}
# Flat set of every curated local tag, for quick membership tests.
_CURATED_LOCAL_SET = {m for fam in CURATED_LOCAL.values() for m in fam}
# OpenAI-compatible chat endpoints. Grok (xAI) and Gemini both expose one, so they
# reuse the same streamer; only Anthropic needs its own Messages API.
OPENAI_COMPAT_BASE: dict[str, str] = {
    "openai": "https://api.openai.com/v1",
    "grok": "https://api.x.ai/v1",
    "gemini": "https://generativelanguage.googleapis.com/v1beta/openai",
}
# Gemini's image (Imagen + native gemini-*-image) and video (Veo) aren't
# OpenAI-compatible — they use the native Generative Language API.
GEMINI_NATIVE_BASE: str = "https://generativelanguage.googleapis.com/v1beta"

# ---- Generated media (image / speech / video) -------------------------------
# Some cloud models don't do chat — they generate an image, audio, or video. We
# route those to the right endpoint and save the bytes to a temp file served by
# /api/media/<id>, so the UI can render <img>/<audio>/<video> without giant
# base64 blobs in the stream. Files live in the OS temp dir (cleared on reboot).
MEDIA_DIR: str = os.path.join(tempfile.gettempdir(), "lls_generated_media")
os.makedirs(MEDIA_DIR, exist_ok=True)
_MEDIA_REGISTRY: dict[str, str] = {}  # served-name -> absolute path (gates traversal)
VIDEO_POLL_SECS: int = int(os.environ.get("LLS_VIDEO_POLL", "5"))      # poll cadence
VIDEO_MAX_WAIT_SECS: int = int(os.environ.get("LLS_VIDEO_MAX_WAIT", "360"))  # give up after


def _save_media(data: bytes, ext: str) -> str:
    """Persist generated bytes to a temp file and return its /api/media/ URL."""
    name = uuid.uuid4().hex + ext
    path = os.path.join(MEDIA_DIR, name)
    with open(path, "wb") as f:
        f.write(data)
    _MEDIA_REGISTRY[name] = path
    return "/api/media/" + name


def _cloud_modality(model: str) -> str:
    """Classify a cloud model by what it OUTPUTS, so non-chat models route to the
    right endpoint instead of erroring as a chat completion. Precise substring
    checks, video before image ('grok-imagine-video' contains neither 'video' issue):
      video → image → audio(TTS) → unsupported(needs input we don't have) → chat."""
    m = (model or "").lower()
    if "video" in m or "sora" in m or "veo" in m:
        return "video"
    if "image" in m or "dall-e" in m or "dalle" in m or "imagen" in m:
        return "image"
    if "tts" in m or "text-to-speech" in m:
        return "audio"
    if any(k in m for k in ("transcribe", "whisper", "realtime", "moderation", "embedding")):
        return "unsupported"
    return "chat"


def _list_provider_models(provider: str, key: str) -> list[str]:
    """Validate a key by listing the account's chat models. Raises on a bad key."""
    if provider == "anthropic":
        r = requests.get(
            "https://api.anthropic.com/v1/models",
            headers={"x-api-key": key, "anthropic-version": "2023-06-01"}, timeout=15,
        )
        r.raise_for_status()
        return [m.get("id", "") for m in r.json().get("data", []) if m.get("id")]
    if provider == "gemini":
        # Native list is the reliable source; stream via the OpenAI-compat endpoint.
        r = requests.get(
            "https://generativelanguage.googleapis.com/v1beta/models",
            params={"key": key}, timeout=15,
        )
        r.raise_for_status()
        out = []
        for m in r.json().get("models", []):
            name = (m.get("name", "") or "").split("/")[-1]
            methods = m.get("supportedGenerationMethods", [])
            if name.startswith("gemini") and (not methods or "generateContent" in methods):
                out.append(name)
        return sorted(set(out))
    base = OPENAI_COMPAT_BASE.get(provider)
    if base:
        r = requests.get(base + "/models", headers={"Authorization": f"Bearer {key}"}, timeout=15)
        r.raise_for_status()
        ids = [m.get("id", "") for m in r.json().get("data", []) if m.get("id")]
        if provider == "openai":
            # Chat-capable only (skip embeddings/tts/whisper/moderation/etc.).
            return sorted(i for i in ids if i.startswith(("gpt", "o1", "o3", "o4")))
        if provider == "grok":
            return sorted(i for i in ids if "grok" in i)
        return sorted(ids)
    return []


def _brave_search(args: dict[str, Any]) -> str:
    query = str(args.get("query") or "").strip()
    if not query:
        return "No search query provided."
    try:
        resp = requests.get(
            "https://api.search.brave.com/res/v1/web/search",
            params={"q": query, "count": 6},
            headers={"Accept": "application/json", "X-Subscription-Token": BRAVE_KEY},
            timeout=15,
        )
        resp.raise_for_status()
        results = ((resp.json().get("web") or {}).get("results") or [])[:6]
    except requests.RequestException as exc:
        return f"Brave search failed: {exc}"
    if not results:
        return f"No web results for '{query}'."
    lines = [
        f"- {r.get('title', '')}\n  {r.get('url', '')}\n  {(r.get('description') or '')[:200]}"
        for r in results
    ]
    return f"Web results for '{query}':\n" + "\n".join(lines)


def _get_weather(args: dict[str, Any]) -> str:
    location = str(args.get("location") or "").strip()
    if not location:
        return "No location provided."
    try:
        resp = requests.get(
            "https://api.openweathermap.org/data/2.5/weather",
            params={"q": location, "appid": WEATHER_KEY, "units": "imperial"},
            timeout=15,
        )
        if resp.status_code == 404:
            return f"Location not found: {location}"
        resp.raise_for_status()
        d = resp.json()
    except requests.RequestException as exc:
        return f"Weather lookup failed: {exc}"
    main = d.get("main") or {}
    desc = (d.get("weather") or [{}])[0].get("description", "")
    wind = (d.get("wind") or {}).get("speed")
    return (
        f"Current weather in {d.get('name', location)}: {desc}, "
        f"{main.get('temp')}°F (feels like {main.get('feels_like')}°F), "
        f"humidity {main.get('humidity')}%, wind {wind} mph."
    )


# Native (built-in) tools the model can call, alongside the MCP tools.
NATIVE_TOOLS: dict[str, Any] = {}
if BRAVE_KEY:
    NATIVE_TOOLS["brave_search"] = _brave_search
if WEATHER_KEY:
    NATIVE_TOOLS["get_weather"] = _get_weather


def _native_tool_schemas() -> list[dict[str, Any]]:
    schemas: list[dict[str, Any]] = []
    if "brave_search" in NATIVE_TOOLS:
        schemas.append({"type": "function", "function": {
            "name": "brave_search",
            "description": "Search the live web with Brave. Returns the top results "
                           "(title, URL, snippet). Use for current events, facts, prices, "
                           "or anything not in the user's own data.",
            "parameters": {"type": "object", "properties": {
                "query": {"type": "string", "description": "the web search query"}},
                "required": ["query"]},
        }})
    if "get_weather" in NATIVE_TOOLS:
        schemas.append({"type": "function", "function": {
            "name": "get_weather",
            "description": "Get the current weather for a place via OpenWeather.",
            "parameters": {"type": "object", "properties": {
                "location": {"type": "string", "description": "city, e.g. 'Williamsburg, VA'"}},
                "required": ["location"]},
        }})
    # Agentic file CRUD (Phase 6) — only offered when the user has turned on
    # "Agent edits". Every change is journaled and undoable (see _crud_* below).
    if AGENT_WRITES_ENABLED:
        schemas += _CRUD_TOOL_SCHEMAS
    return schemas


# A drafting tool, not a sending one. The model uses it to PROPOSE a text; the app
# shows the user a Send/Cancel card and only sends if they tap Send. The real
# send tool stays filtered out of the model's read-only tool set, so the model can
# never send a message on its own.
DRAFT_IMESSAGE_SCHEMA: dict[str, Any] = {
    "type": "function",
    "function": {
        "name": "draft_imessage",
        "description": (
            "Prepare an iMessage/SMS for the user to review. This does NOT send the "
            "message — it shows the user a Send/Cancel card, and the text is sent only "
            "if they tap Send. Use this whenever the user asks you to text or message "
            "someone. Call it once, then briefly tell the user you've prepared the draft."
        ),
        "parameters": {"type": "object", "properties": {
            "to": {"type": "string", "description": "recipient: a phone number, email, or contact name"},
            "text": {"type": "string", "description": "the message body to send"}},
            "required": ["to", "text"]},
    },
}


# ---- Shared live document (one open .md; the "pen" gates who may write) ------
# Single source of truth, server-side. All models READ it (injected into context);
# only the pen-holder's chat request carries can_edit_doc=true, so only that model
# is handed the edit_document tool. The pen is the frontend's selected pane.
OPEN_DOC: dict[str, Any] = {"path": "", "content": "", "kind": "", "editable": False}
_MD_EXTS = (".md", ".markdown")
_IMG_EXTS = (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg", ".ico")
_VID_EXTS = (".mp4", ".webm", ".mov", ".m4v", ".ogg")


def _doc_kind(path: str, is_text: bool) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in _MD_EXTS:
        return "markdown"
    if ext in _IMG_EXTS:
        return "image"
    if ext == ".pdf":
        return "pdf"
    if ext in _VID_EXTS:
        return "video"
    return "text" if is_text else "binary"


def _open_doc_path(raw: str) -> Response:
    """Open ANY file the user picked into the live viewer. Text files are read and
    become editable (the pen); binary files (image/pdf/video/…) are served raw for
    preview and are not text-editable. The user chooses the file explicitly (native
    picker), so any readable file is allowed."""
    if not raw:
        return jsonify({"ok": False, "error": "No file selected."}), 400
    path = os.path.realpath(os.path.expanduser(raw))
    if not os.path.isfile(path):
        return jsonify({"ok": False, "error": "File not found."}), 404
    content, is_text = "", False
    try:
        with open(path, encoding="utf-8") as f:
            content = f.read()
        is_text = True
    except (UnicodeDecodeError, OSError):
        is_text = False
    kind = _doc_kind(path, is_text)
    OPEN_DOC.update({"path": path, "content": content if is_text else "", "kind": kind, "editable": is_text})
    return jsonify({
        "ok": True, "path": path, "name": os.path.basename(path),
        "content": content if is_text else "", "kind": kind,
        "editable": is_text, "size": os.path.getsize(path),
    })


EDIT_DOCUMENT_SCHEMA: dict[str, Any] = {
    "type": "function",
    "function": {
        "name": "edit_document",
        "description": (
            "Edit the shared document open in the live viewer. You currently hold the "
            "pen, so your change is saved immediately. Give `find` (exact text already "
            "in the document) and `replace` to change a passage; OR omit `find` to "
            "append `replace` to the end. Keep edits surgical and preserve the rest."
        ),
        "parameters": {"type": "object", "properties": {
            "find": {"type": "string", "description": "exact existing text to replace; omit or leave empty to append"},
            "replace": {"type": "string", "description": "the new text"}},
            "required": ["replace"]},
    },
}


def _apply_doc_edit(find: str, replace: str) -> str:
    """Apply a find/replace (or append) to the open doc. Scoped to OPEN_DOC only."""
    path = OPEN_DOC["path"]
    if not path:
        return "No document is open."
    try:
        with open(path, encoding="utf-8") as f:
            content = f.read()
    except OSError as exc:
        return f"Could not read the document: {exc}"
    if find:
        if find not in content:
            return "The `find` text was not found in the document; no change made."
        new = content.replace(find, replace, 1)
    else:
        new = content.rstrip("\n") + "\n\n" + replace + "\n"
    try:
        with open(path, "w", encoding="utf-8") as f:
            f.write(new)
    except OSError as exc:
        return f"Could not write the document: {exc}"
    OPEN_DOC["content"] = new
    return "Document updated."


def _dispatch_tool(name: str, args: dict[str, Any]) -> str:
    """Route a tool call to a native tool (Brave/weather) or the MCP bridge."""
    fn = NATIVE_TOOLS.get(name)
    return fn(args) if fn else MCP.call(name, args)

# Injected as the system message on every chat so the model knows the app can
# render artifacts it writes, and what the user can do with them. Overridable
# via the LLS_SYSTEM_PROMPT env var.
_DEFAULT_SYSTEM_PROMPT: str = (
    "You are a helpful assistant running fully offline inside Local LLM Studio, a "
    "local chat app. The app renders three kinds of fenced code blocks into live, "
    "interactive results, and the user can Save them to a gallery, Print them, and "
    "Download them at any time. Reach for these whenever a visual would make your "
    "answer clearer — you may create them proactively, not only when asked:\n\n"
    "- ```mermaid — diagrams: flowchart, sequence, class, ER, state, gantt, "
    "mindmap, pie. Best for processes, architectures, timelines, and relationships.\n"
    "- ```svg — a complete <svg>...</svg> for custom vector graphics, simple charts, "
    "or icons.\n"
    "- ```chart — a Chart.js config as JSON, e.g. "
    '{"type":"bar","data":{"labels":["A","B"],"datasets":[{"label":"X","data":[3,7]}]}}. '
    "type can be bar, line, pie, doughnut, radar, scatter, or polarArea. Use this to "
    "plot actual data and numbers as a graph or chart.\n"
    "- ```html — a self-contained HTML document (inline CSS and JS) shown in a "
    "sandboxed iframe, for interactive widgets, styled tables, or small demos.\n\n"
    "Because everything runs offline, keep artifacts self-contained: inline all CSS "
    "and JS, and do NOT reference external URLs, CDNs, web fonts, or remote images. "
    "Keep the syntax valid, put only the artifact inside its fence, and add a brief "
    "explanation outside it. Use normal markdown prose for everything else, and feel "
    "free to remind the user they can save, print, or download any artifact you make."
)
def _load_system_prompt() -> str:
    """Base system prompt / global instructions: an explicit env override, else the
    editable instructions.md beside the app, else the built-in default."""
    override = os.environ.get("LLS_SYSTEM_PROMPT")
    if override:
        return override
    path = os.environ.get(
        "LLS_SYSTEM_PROMPT_FILE",
        os.path.join(os.path.dirname(os.path.abspath(__file__)), "instructions.md"),
    )
    try:
        with open(path, encoding="utf-8") as f:
            text = f.read().strip()
        if text:
            return text
    except OSError:
        pass
    return _DEFAULT_SYSTEM_PROMPT


SYSTEM_PROMPT: str = _load_system_prompt()

# Set by the supervisor when it spawns this worker; lets the UI's End / Settings
# controls reach the lifecycle manager. Empty when run standalone (no controls).
CONTROL_URL: str = os.environ.get("LLS_CONTROL_URL", "")


def _human_size(num_bytes: int) -> str:
    """Render a byte count as a compact human-readable string (e.g. '3.8 GB')."""
    size = float(num_bytes)
    for unit in ("B", "KB", "MB", "GB", "TB"):
        if size < 1024.0 or unit == "TB":
            return f"{size:.0f} {unit}" if unit == "B" else f"{size:.1f} {unit}"
        size /= 1024.0
    return f"{size:.1f} TB"


FAVICON_SVG = (
    "<svg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 32 32'>"
    "<rect width='32' height='32' rx='7' fill='#0D1117'/>"
    "<path d='M9 8h14a3 3 0 0 1 3 3v7a3 3 0 0 1-3 3h-7l-5 4v-4h-2a3 3 0 0 1-3-3v-7a3 3 0 0 1 3-3z' "
    "fill='#D4A853'/></svg>"
)


def _asset_version() -> str:
    """Max mtime of the front-end assets, appended as ?v= so the browser always
    fetches fresh JS/CSS the moment a file changes — no hard-reload needed."""
    sdir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "static")
    files = ["app.js", "style.css", "markdown.js", os.path.join("vendor", "highlight-mini.js")]
    try:
        return str(int(max(os.path.getmtime(os.path.join(sdir, f)) for f in files)))
    except OSError:
        return "1"


@app.get("/")
def index() -> str:
    """Serve the single-page chat UI."""
    return render_template("index.html", control_url=CONTROL_URL, asset_v=_asset_version())


@app.get("/favicon.ico")
def favicon() -> Response:
    """Serve the inline brand mark so /favicon.ico never 404s, even on a direct hit."""
    return Response(FAVICON_SVG, mimetype="image/svg+xml")


@app.get("/local")
def local_file() -> Response:
    """Serve a local file so the UI can embed it (image/video/pdf/etc.). Restricted
    to the user's home directory with traversal blocked — never serves system files."""
    raw = request.args.get("path", "")
    if not raw:
        return jsonify({"error": "path required"}), 400
    path = os.path.realpath(os.path.expanduser(raw))
    home = os.path.realpath(os.path.expanduser("~"))
    if not (path == home or path.startswith(home + os.sep)) or not os.path.isfile(path):
        return jsonify({"error": "not found or outside the allowed area"}), 404
    return send_file(path)


@app.get("/api/media/<name>")
def api_media(name: str) -> Response:
    """Serve a file we generated (image/audio/video). Resolves from the in-memory
    registry (ephemeral) or the persistent saved-media dir; the name is validated so
    there's no path traversal and we never serve arbitrary files."""
    path = _MEDIA_REGISTRY.get(name)
    if not path and re.fullmatch(r"[A-Za-z0-9._-]+", name or ""):
        cand = os.path.join(ARTIFACT_MEDIA_DIR, name)
        if os.path.isfile(cand):
            path = cand
    if not path or not os.path.isfile(path):
        return jsonify({"error": "not found"}), 404
    return send_file(path)


@app.post("/api/artifacts/media")
def save_media_artifact() -> Response:
    """Persist a generated image/audio/video and save an 'embed' artifact pointing at
    it, so it shows up in the gallery and survives a restart. Body: {url, kind, title?}.
    `url` is either our /api/media/<name> or a remote provider URL (downloaded here)."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    url = str(data.get("url", "")).strip()
    kind = str(data.get("kind", "image"))
    title = (str(data.get("title") or "").strip() or f"Generated {kind}")[:120]
    default_ext = {"image": ".png", "audio": ".mp3", "video": ".mp4"}.get(kind, ".bin")
    blob, ext = None, default_ext
    if url.startswith("/api/media/"):
        name = url.rsplit("/", 1)[-1]
        src = _MEDIA_REGISTRY.get(name)
        if not src and re.fullmatch(r"[A-Za-z0-9._-]+", name):
            cand = os.path.join(ARTIFACT_MEDIA_DIR, name)
            src = cand if os.path.isfile(cand) else None
        if not src or not os.path.isfile(src):
            return jsonify({"ok": False, "error": "That media is no longer available to save."}), 404
        ext = os.path.splitext(src)[1] or default_ext
        with open(src, "rb") as f:
            blob = f.read()
    elif url.startswith("http://") or url.startswith("https://"):
        try:
            r = requests.get(url, timeout=60)
            if r.status_code >= 400:
                return jsonify({"ok": False, "error": "Could not fetch the media to save."}), 502
            blob = r.content
        except requests.RequestException as exc:
            return jsonify({"ok": False, "error": f"Fetch failed: {exc}"}), 502
    else:
        return jsonify({"ok": False, "error": "Unsupported media URL."}), 400

    os.makedirs(ARTIFACT_MEDIA_DIR, exist_ok=True)
    name = uuid.uuid4().hex + ext
    with open(os.path.join(ARTIFACT_MEDIA_DIR, name), "wb") as f:
        f.write(blob)
    aid = datetime.now().strftime("%Y%m%d-%H%M%S-") + secrets.token_hex(3)
    record = {
        "id": aid, "type": "embed", "title": title,
        "source": "/api/media/" + name,
        "created": datetime.now(timezone.utc).isoformat(),
    }
    os.makedirs(ARTIFACTS_DIR, exist_ok=True)
    with open(_artifact_path(aid), "w", encoding="utf-8") as f:
        json.dump(record, f, ensure_ascii=False)
    return jsonify({"ok": True, **{k: record[k] for k in ("id", "type", "title", "created")}})


# ---- Voice mode (#49): cloud STT (Whisper) + TTS, via the OpenAI key ----------
VOICE_STT_MODEL: str = os.environ.get("LLS_VOICE_STT_MODEL", "whisper-1")
VOICE_TTS_MODEL: str = os.environ.get("LLS_VOICE_TTS_MODEL", "gpt-4o-mini-tts")
VOICE_TTS_VOICE: str = os.environ.get("LLS_VOICE_TTS_VOICE", "alloy")


@app.post("/api/transcribe")
def api_transcribe() -> Response:
    """Speech-to-text for voice mode: forward the browser's mic recording to OpenAI
    Whisper and return the transcript. Requires the OpenAI key."""
    key = PROVIDER_KEYS.get("openai", "")
    if not key:
        return jsonify({"ok": False, "error": "Connect an OpenAI key in Settings → Cloud models to use voice."}), 400
    f = request.files.get("audio")
    if not f or not f.filename:
        return jsonify({"ok": False, "error": "No audio was uploaded."}), 400
    try:
        r = requests.post(
            OPENAI_COMPAT_BASE["openai"] + "/audio/transcriptions",
            headers={"Authorization": f"Bearer {key}"},
            files={"file": (f.filename or "audio.webm", f.read(), f.mimetype or "audio/webm")},
            data={"model": request.form.get("model") or VOICE_STT_MODEL}, timeout=120,
        )
    except requests.RequestException as exc:
        return jsonify({"ok": False, "error": f"Transcription request failed: {exc}"}), 502
    if r.status_code >= 400:
        return jsonify({"ok": False, "error": _cloud_err(r, "Transcription")}), 502
    try:
        return jsonify({"ok": True, "text": (r.json().get("text") or "").strip()})
    except ValueError:
        return jsonify({"ok": True, "text": (r.text or "").strip()})


@app.post("/api/speak")
def api_speak() -> Response:
    """Text-to-speech for voice mode: return mp3 bytes for immediate playback in the
    browser. Requires the OpenAI key."""
    key = PROVIDER_KEYS.get("openai", "")
    if not key:
        return jsonify({"ok": False, "error": "Connect an OpenAI key to use voice."}), 400
    data = request.get_json(silent=True) or {}
    text = str(data.get("text", "")).strip()
    if not text:
        return jsonify({"ok": False, "error": "No text to speak."}), 400
    try:
        r = requests.post(
            OPENAI_COMPAT_BASE["openai"] + "/audio/speech",
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
            json={"model": VOICE_TTS_MODEL, "input": text[:4000],
                  "voice": str(data.get("voice") or VOICE_TTS_VOICE), "response_format": "mp3"},
            timeout=120,
        )
    except requests.RequestException as exc:
        return jsonify({"ok": False, "error": f"Speech request failed: {exc}"}), 502
    if r.status_code >= 400:
        return jsonify({"ok": False, "error": _cloud_err(r, "Speech")}), 502
    return Response(r.content, mimetype="audio/mpeg")


@app.get("/api/graph")
def api_graph() -> Response:
    """The local Graphify knowledge graph (node-link JSON), trimmed to what the
    Second-Brain view needs. Empty + available:false if it hasn't been built."""
    try:
        with open(RAG_GRAPH, encoding="utf-8") as f:
            g = json.load(f)
    except (OSError, ValueError):
        return jsonify({"available": False, "nodes": [], "links": []})
    nodes = [
        {"id": n.get("id"), "label": n.get("label") or n.get("id"),
         "community": n.get("community", 0), "type": n.get("file_type", ""),
         # The source note, if it resolves to a real .md — lets the UI open it on click.
         "file": (n.get("source_file", "") if _resolve_note_path(n.get("source_file", "")) else "")}
        for n in g.get("nodes", []) if n.get("id")
    ]
    valid = {n["id"] for n in nodes}
    links = [
        {"source": e.get("source"), "target": e.get("target"),
         "relation": e.get("relation", ""), "confidence": e.get("confidence_score", 1)}
        for e in g.get("links", [])
        if e.get("source") in valid and e.get("target") in valid
    ]
    return jsonify({"available": True, "nodes": nodes, "links": links,
                    "commit": g.get("built_at_commit", "")})


@app.get("/api/note")
def api_note() -> Response:
    """Return a graph node's source note as markdown, for the click-to-preview
    overlay. `file` is the node's source_file; resolution is traversal-guarded."""
    path = _resolve_note_path(request.args.get("file", ""))
    if not path:
        return jsonify({"ok": False, "error": "Note not found."}), 404
    try:
        with open(path, encoding="utf-8") as f:
            md = f.read()
    except OSError as exc:
        return jsonify({"ok": False, "error": f"Could not read note: {exc}"}), 500
    return jsonify({"ok": True, "name": os.path.basename(path), "markdown": md})


@app.get("/api/secondbrain/health")
def api_secondbrain_health() -> Response:
    """Index health for the Second-Brain panel: embed model, note count, last reindex."""
    info: dict[str, Any] = {"available": False, "model": RAG_MODEL, "notes": 0, "reindexed": None}
    try:
        st = os.stat(RAG_INDEX)
        with open(RAG_INDEX, encoding="utf-8") as f:
            idx = json.load(f)
        info["available"] = True
        info["notes"] = len(idx)
        info["reindexed"] = datetime.fromtimestamp(st.st_mtime).astimezone().isoformat()
    except (OSError, ValueError):
        pass
    return jsonify(info)


_TEXT_EXTS = frozenset({
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".json", ".log", ".rtf",
    ".py", ".js", ".ts", ".html", ".css", ".yaml", ".yml", ".xml",
})


@app.post("/api/extract")
def extract_text() -> Response:
    """Extract text from an uploaded document (pdf / pptx / docx / text-like) so the
    model can read it. Truncated to fit the context window. Images are handled
    client-side as base64 (sent to a vision model), not here."""
    import io

    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": "no file uploaded"}), 400
    name = f.filename
    ext = os.path.splitext(name)[1].lower()
    data = f.read()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(data))
            text = "\n".join((page.extract_text() or "") for page in reader.pages)
        elif ext == ".pptx":
            from pptx import Presentation
            parts: list[str] = []
            for i, slide in enumerate(Presentation(io.BytesIO(data)).slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        parts.append(shape.text_frame.text)
            text = "\n".join(parts)
        elif ext == ".docx":
            from docx import Document
            text = "\n".join(p.text for p in Document(io.BytesIO(data)).paragraphs)
        elif ext in _TEXT_EXTS:
            text = data.decode("utf-8", "replace")
        else:
            return jsonify({
                "error": f"Unsupported document type: {ext or '(none)'}. "
                         "Images attach directly; for documents use pdf, pptx, docx, txt, md, or csv."
            }), 415
    except Exception as exc:  # malformed/encrypted file, etc.
        return jsonify({"error": f"Could not read {name}: {exc}"}), 422

    text = text.strip()
    truncated = len(text) > EXTRACT_MAX_CHARS
    if truncated:
        text = text[:EXTRACT_MAX_CHARS]
    return jsonify({"filename": name, "text": text, "chars": len(text), "truncated": truncated})


@app.get("/api/health")
def health() -> Response:
    """Quick reachability check of the Ollama server."""
    try:
        resp = requests.get(f"{OLLAMA_HOST}/api/tags", timeout=REQUEST_TIMEOUT)
        resp.raise_for_status()
        return jsonify({"ok": True})
    except requests.RequestException:
        return jsonify({"ok": False})


@app.get("/api/ps")
def ps() -> Response:
    """Loaded models + real GPU/VRAM offload, from Ollama's /api/ps. Used by the
    Telemetry HUD. (We surface only metrics we can read truthfully — GPU%/temp/power
    need privileged powermetrics, so the HUD omits them rather than fake them.)"""
    try:
        resp = requests.get(f"{OLLAMA_HOST}/api/ps", timeout=REQUEST_TIMEOUT)
        resp.raise_for_status()
        data = resp.json()
    except (requests.RequestException, ValueError):
        return jsonify({"models": []})
    out = []
    for m in data.get("models", []):
        size = int(m.get("size", 0) or 0)
        vram = int(m.get("size_vram", 0) or 0)
        out.append({
            "name": m.get("name") or m.get("model", "?"),
            "size": size,
            "size_vram": vram,
            "size_h": _human_size(size),
            "vram_h": _human_size(vram),
            "gpu_pct": round(100 * vram / size) if size else 0,
        })
    return jsonify({"models": out})


@app.get("/api/models")
def models() -> Response:
    """The curated local models (grouped by provider), each flagged 'installed'
    or downloadable. Only these are offered — any other installed model is hidden.
    Selecting a not-installed one triggers /api/pull in the UI."""
    installed: dict[str, str] = {}  # name -> human size
    try:
        resp = requests.get(f"{OLLAMA_HOST}/api/tags", timeout=REQUEST_TIMEOUT)
        resp.raise_for_status()
        for m in resp.json().get("models", []):
            installed[m.get("name", "")] = _human_size(int(m.get("size", 0)))
    except requests.RequestException:
        pass  # Ollama unreachable → all show as downloadable; pull reports the real error
    out: list[dict[str, Any]] = []
    for group, tags in CURATED_LOCAL.items():
        for tag in tags:
            out.append({
                "name": tag,
                "group": group,
                "installed": tag in installed,
                "size": installed.get(tag, ""),
            })
    return jsonify({"models": out})


@app.post("/api/pull")
def api_pull() -> Response:
    """Download a curated local model via Ollama, streaming progress to the UI as
    NDJSON ({status, total, completed, ...}). Whitelisted to the curated set so it
    can never be told to pull an arbitrary model."""
    data = request.get_json(silent=True) or {}
    model = str(data.get("model") or "").strip()
    if model not in _CURATED_LOCAL_SET:
        return jsonify({"error": "That model isn't in the studio's model list."}), 400

    def generate() -> Iterator[bytes]:
        try:
            with requests.post(
                f"{OLLAMA_HOST}/api/pull", json={"model": model, "stream": True},
                stream=True, timeout=None,
            ) as up:
                if up.status_code >= 400:
                    yield (json.dumps({"status": "error",
                                       "error": f"Ollama returned {up.status_code}."}) + "\n").encode("utf-8")
                    return
                for line in up.iter_lines():
                    if line:
                        yield line + b"\n"
        except requests.RequestException as exc:
            yield (json.dumps({"status": "error", "error": str(exc)}) + "\n").encode("utf-8")

    return Response(generate(), mimetype="application/x-ndjson")


@app.get("/api/capabilities")
def capabilities() -> Response:
    """Report whether a model can stream reasoning ('thinking'), via /api/show.

    The UI uses this to decide automatically — no toggle — whether to request
    thinking for the selected model. Always 200s with thinking:false on any
    failure so the chat UI degrades gracefully.
    """
    name = request.args.get("model", "").strip()
    if not name:
        return jsonify({"error": "Query param 'model' is required."}), 400
    caps: list[Any] = []
    ctx = 0
    try:
        resp = requests.post(
            f"{OLLAMA_HOST}/api/show", json={"model": name}, timeout=REQUEST_TIMEOUT
        )
        resp.raise_for_status()
        body = resp.json()
        caps = body.get("capabilities") or []
        mi = body.get("model_info") or {}
        ctx = next((v for k, v in mi.items()
                    if k.endswith("context_length") and isinstance(v, int)), 0)
    except (requests.RequestException, ValueError):
        pass
    if ctx:
        _CTX_CACHE[name] = ctx
    return jsonify({
        "model": name,
        "thinking": "thinking" in caps,
        "vision": "vision" in caps,
        "capabilities": caps,
        "context_length": ctx,
        "num_ctx": min(ctx, NUM_CTX_CAP) if ctx else NUM_CTX_CAP,
    })


_TOOLS_CACHE: dict[str, bool] = {}


def _model_supports_tools(model: str) -> bool:
    """Whether a local model can tool-call (Ollama 'tools' capability), cached.
    Models like gemma3 lack it; sending tools makes Ollama 400 with 'does not
    support tools', so the chat path must check before using the tool loop."""
    if model in _TOOLS_CACHE:
        return _TOOLS_CACHE[model]
    supported = False
    try:
        resp = requests.post(
            f"{OLLAMA_HOST}/api/show", json={"model": model}, timeout=REQUEST_TIMEOUT
        )
        resp.raise_for_status()
        supported = "tools" in (resp.json().get("capabilities") or [])
    except (requests.RequestException, ValueError):
        supported = False
    _TOOLS_CACHE[model] = supported
    return supported


def _effective_num_ctx(model: str) -> int:
    """The context window we ask Ollama to use for a model: its trained max, capped."""
    ctx = _CTX_CACHE.get(model, 0)
    if not ctx:
        try:
            resp = requests.post(
                f"{OLLAMA_HOST}/api/show", json={"model": model}, timeout=REQUEST_TIMEOUT
            )
            resp.raise_for_status()
            mi = resp.json().get("model_info") or {}
            ctx = next((v for k, v in mi.items()
                        if k.endswith("context_length") and isinstance(v, int)), 0)
            if ctx:
                _CTX_CACHE[model] = ctx
        except (requests.RequestException, ValueError):
            ctx = 0
    return min(ctx, NUM_CTX_CAP) if ctx else NUM_CTX_CAP


# ---- Saved artifacts (diagrams) -------------------------------------------
def _artifact_path(aid: str) -> str:
    return os.path.join(ARTIFACTS_DIR, f"{aid}.json")


@app.get("/api/artifacts")
def list_artifacts() -> Response:
    """List saved artifacts (metadata only), newest first."""
    items: list[dict[str, Any]] = []
    if os.path.isdir(ARTIFACTS_DIR):
        for fn in os.listdir(ARTIFACTS_DIR):
            if not fn.endswith(".json"):
                continue
            try:
                with open(os.path.join(ARTIFACTS_DIR, fn), encoding="utf-8") as f:
                    rec = json.load(f)
                items.append(
                    {k: rec[k] for k in ("id", "type", "title", "created")}
                )
            except (OSError, ValueError, KeyError):
                continue
    items.sort(key=lambda x: x["created"], reverse=True)
    return jsonify({"artifacts": items})


@app.post("/api/artifacts")
def save_artifact() -> Response:
    """Persist a diagram/artifact to disk. Body: {type, source, title?}."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    atype = data.get("type")
    source = data.get("source")
    title = (str(data.get("title") or "").strip() or "Untitled")[:120]
    if atype not in ARTIFACT_TYPES:
        return jsonify({"error": "type must be one of mermaid, svg, html"}), 400
    if not isinstance(source, str) or not source.strip():
        return jsonify({"error": "source is required"}), 400
    if len(source.encode("utf-8")) > MAX_ARTIFACT_BYTES:
        return jsonify({"error": "source too large"}), 413

    aid = datetime.now().strftime("%Y%m%d-%H%M%S-") + secrets.token_hex(3)
    record = {
        "id": aid,
        "type": atype,
        "title": title,
        "source": source,
        "created": datetime.now(timezone.utc).isoformat(),
    }
    os.makedirs(ARTIFACTS_DIR, exist_ok=True)
    with open(_artifact_path(aid), "w", encoding="utf-8") as f:
        json.dump(record, f, ensure_ascii=False)
    return jsonify({k: record[k] for k in ("id", "type", "title", "created")})


@app.get("/api/artifacts/<aid>")
def get_artifact(aid: str) -> Response:
    """Return one saved artifact in full (including source)."""
    if not ARTIFACT_ID_RE.match(aid):
        return jsonify({"error": "bad id"}), 400
    path = _artifact_path(aid)
    if not os.path.exists(path):
        return jsonify({"error": "not found"}), 404
    with open(path, encoding="utf-8") as f:
        return jsonify(json.load(f))


@app.delete("/api/artifacts/<aid>")
def delete_artifact(aid: str) -> Response:
    """Delete one saved artifact (user-initiated, from the gallery)."""
    if not ARTIFACT_ID_RE.match(aid):
        return jsonify({"error": "bad id"}), 400
    path = _artifact_path(aid)
    if os.path.exists(path):
        os.remove(path)
    return jsonify({"deleted": aid})


# ---- Second-brain grounding ------------------------------------------------
def _read_note(path: str) -> tuple[str, str]:
    """Return (title, body) for a vault note, stripping YAML frontmatter."""
    with open(path, encoding="utf-8") as f:
        text = f.read()
    m = re.match(r"^---\n(.*?)\n---\n?(.*)$", text, re.S)
    if m:
        tm = re.search(r"^title:\s*(.+)$", m.group(1), re.M)
        title = tm.group(1).strip() if tm else os.path.splitext(os.path.basename(path))[0]
        return title, m.group(2).strip()
    return os.path.splitext(os.path.basename(path))[0], text.strip()


def _retrieve_context(query: str, exclude_private: bool = False) -> tuple[str, list[dict[str, Any]]]:
    """Ask the local skippy-rag index for notes relevant to `query`. Returns a
    (context_block, cited_notes) pair; ("", []) when nothing clears the bar or the
    index/Ollama is unavailable (so the chat just proceeds ungrounded).

    The index includes the vault's `private/` notes (message logs, the ward directory).
    For a **cloud** model, pass exclude_private=True so that PII is NEVER injected into a
    prompt that leaves the machine — a structural guard that doesn't depend on the model
    choosing not to quote it. Local/offline models retrieve private notes normally."""
    if not (os.path.exists(RAG_PYTHON) and os.path.exists(RAG_QUERY)):
        return "", []
    try:
        proc = subprocess.run(
            [RAG_PYTHON, RAG_QUERY, query, "--top-k", str(RAG_TOP_K)],
            capture_output=True, text=True, timeout=RAG_TIMEOUT, cwd=RAG_DIR,
        )
    except (subprocess.SubprocessError, OSError):
        return "", []
    if proc.returncode != 0:
        return "", []

    blocks: list[str] = []
    cited: list[dict[str, Any]] = []
    for line in proc.stdout.splitlines():
        parts = line.split("\t")
        if len(parts) < 2:
            continue
        try:
            score = float(parts[0])
        except ValueError:
            continue
        path = parts[1]
        if score < RAG_MIN_SCORE or not os.path.exists(path):
            continue
        # Cloud guard: never inject vault `private/` PII into a prompt that leaves the machine.
        if exclude_private and os.path.realpath(path).startswith(
            os.path.realpath(os.path.join(VAULT_ROOT, "private")) + os.sep
        ):
            continue
        try:
            title, body = _read_note(path)
        except OSError:
            continue
        blocks.append(f"## {title}\n{body[:RAG_MAX_CHARS]}")
        cited.append({
            "title": title,
            "score": round(score, 3),
            "stem": os.path.splitext(os.path.basename(path))[0],
            "snippet": re.sub(r"\s+", " ", body).strip()[:220],
        })

    if not blocks:
        return "", []
    context = (
        "\n\n# Context from the user's second brain\n"
        "These are the user's own notes, retrieved as possibly relevant to their "
        "message. Ground your answer in them and cite note titles you rely on. If "
        "they are not relevant, ignore them.\n\n" + "\n\n".join(blocks)
    )
    return context, cited


def _nd(obj: dict[str, Any]) -> bytes:
    return (json.dumps(obj) + "\n").encode("utf-8")


def _execute_tool(name: str, args: dict[str, Any]) -> tuple[str, list[bytes]]:
    """Run one tool call, returning (result_text_for_model, ndjson_events_for_ui).

    The single place every tool loop — local (Ollama), Anthropic, and
    OpenAI-compatible — runs a tool, so the special-cased tools behave identically
    across providers:
      • draft_imessage never sends; it emits a Send/Cancel card for the user.
      • edit_document writes the open shared doc (the caller only reaches here when
        it holds the pen) and emits a doc_event so the UI refreshes.
    Everything else dispatches to a native (Brave/weather) or read-only MCP tool.
    """
    if name == "draft_imessage":
        to = str(args.get("to") or "").strip()
        text = str(args.get("text") or "").strip()
        result = (
            "A confirmation card with this draft was shown to the user. The message "
            "will be sent ONLY if they tap Send — you cannot send it yourself. Do not "
            "call draft_imessage again; just briefly confirm you've prepared the draft "
            "for their review."
        )
        return result, [_nd({"confirm_send": {"to": to, "text": text}})]
    if name == "edit_document":
        result = _apply_doc_edit(str(args.get("find") or ""), str(args.get("replace") or ""))
        return result, [_nd({"doc_event": {"updated": result == "Document updated.", "note": result}})]
    # Agentic file change with approval required → stage it and show an Approve/
    # Reject card instead of applying. The model can never apply it itself.
    if name in _CRUD_TOOL_NAMES and AGENT_WRITES_ENABLED and AGENT_APPROVAL_REQUIRED:
        return _stage_change(name, args)
    events = [_nd({"tool_event": {"phase": "call", "name": name, "args": args}})]
    result = _dispatch_tool(name, args)
    events.append(_nd({"tool_event": {"phase": "result", "name": name, "preview": result[:240]}}))
    return result, events


def _anthropic_tool_schemas(tools: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Translate the app's OpenAI-style function schemas into Anthropic's tool shape
    ({name, description, input_schema}). Anthropic is the one provider that doesn't
    take the OpenAI format verbatim."""
    out: list[dict[str, Any]] = []
    for t in tools or []:
        fn = t.get("function") or {}
        if not fn.get("name"):
            continue
        out.append({
            "name": fn["name"],
            "description": fn.get("description", ""),
            "input_schema": fn.get("parameters") or {"type": "object", "properties": {}},
        })
    return out


def _stream_anthropic(
    model: str, messages: list[dict[str, Any]], options: Any, key: str,
    tools: list[dict[str, Any]] | None = None,
) -> Iterator[bytes]:
    """Stream Anthropic's Messages API, translated into the app's NDJSON shape
    ({"message":{"content"|"thinking":…}} chunks + a final {"done":true,…} with token
    counts). When `tools` is given, runs the tool_use loop: stream a turn, run any
    requested tools via _execute_tool, feed the tool_result blocks back, and repeat
    until Claude stops asking for tools (or MAX_TOOL_ROUNDS)."""
    system = ""
    conv: list[dict[str, Any]] = []
    for m in messages:
        role, content = m.get("role"), m.get("content", "")
        if not isinstance(content, str) or not content:
            continue
        if role == "system":
            system += ("\n\n" if system else "") + content
        elif role in ("user", "assistant"):
            conv.append({"role": role, "content": content})
    atools = _anthropic_tool_schemas(tools) if tools else []
    in_tok = out_tok = 0
    for _round in range(MAX_TOOL_ROUNDS if atools else 1):
        body: dict[str, Any] = {"model": model, "messages": conv, "max_tokens": 4096, "stream": True}
        if system:
            body["system"] = system
        if atools:
            body["tools"] = atools
        # Note: newer Claude models reject the `temperature` param (returns 400
        # "temperature is deprecated for this model"), so we omit it and let
        # Anthropic use its default. Don't reintroduce it for the Messages API.
        blocks: dict[int, dict[str, Any]] = {}  # index -> {type, text|(id,name,json)}
        stop_reason = None
        round_in = round_out = 0
        try:
            with requests.post(
                "https://api.anthropic.com/v1/messages", json=body, stream=True, timeout=None,
                headers={"x-api-key": key, "anthropic-version": "2023-06-01", "content-type": "application/json"},
            ) as up:
                if up.status_code >= 400:
                    try:
                        detail = up.json().get("error", {}).get("message", "")
                    except ValueError:
                        detail = (up.text or "")[:300]
                    yield _nd({"error": detail or f"Anthropic returned {up.status_code}.", "done": True})
                    return
                for raw in up.iter_lines():
                    if not raw:
                        continue
                    line = raw.decode("utf-8", "replace").strip()
                    if not line.startswith("data:"):
                        continue
                    try:
                        ev = json.loads(line[5:].strip())
                    except ValueError:
                        continue
                    t = ev.get("type")
                    if t == "message_start":
                        round_in = (ev.get("message", {}).get("usage", {}) or {}).get("input_tokens", round_in)
                    elif t == "content_block_start":
                        cb = ev.get("content_block") or {}
                        idx = ev.get("index", 0)
                        if cb.get("type") == "tool_use":
                            blocks[idx] = {"type": "tool_use", "id": cb.get("id"), "name": cb.get("name"), "json": ""}
                        else:
                            blocks[idx] = {"type": "text", "text": ""}
                    elif t == "content_block_delta":
                        idx = ev.get("index", 0)
                        d = ev.get("delta", {})
                        if d.get("type") == "text_delta" and d.get("text"):
                            yield _nd({"message": {"content": d["text"]}})
                            blocks.setdefault(idx, {"type": "text", "text": ""})["text"] += d["text"]
                        elif d.get("type") == "thinking_delta" and d.get("thinking"):
                            yield _nd({"message": {"thinking": d["thinking"]}})
                        elif d.get("type") == "input_json_delta":
                            blocks.setdefault(idx, {"type": "tool_use", "id": None, "name": None, "json": ""})["json"] += d.get("partial_json", "")
                    elif t == "message_delta":
                        if (ev.get("delta") or {}).get("stop_reason"):
                            stop_reason = ev["delta"]["stop_reason"]
                        round_out = (ev.get("usage", {}) or {}).get("output_tokens", round_out)
                    elif t == "message_stop":
                        break
        except requests.RequestException as exc:
            yield _nd({"error": f"Anthropic request failed: {exc}", "done": True})
            return

        in_tok += round_in
        out_tok += round_out
        tool_uses = [b for b in blocks.values() if b.get("type") == "tool_use" and b.get("name")]
        if not atools or stop_reason != "tool_use" or not tool_uses:
            yield _nd({"done": True, "prompt_eval_count": in_tok, "eval_count": out_tok})
            return

        # Re-create the assistant turn (text + tool_use blocks) for the next request,
        # then run each tool and feed the results back as a user turn.
        assistant_content: list[dict[str, Any]] = []
        for idx in sorted(blocks):
            b = blocks[idx]
            if b["type"] == "text" and b.get("text"):
                assistant_content.append({"type": "text", "text": b["text"]})
            elif b["type"] == "tool_use" and b.get("name"):
                try:
                    inp = json.loads(b["json"]) if b.get("json") else {}
                except ValueError:
                    inp = {}
                assistant_content.append({"type": "tool_use", "id": b["id"], "name": b["name"], "input": inp})
        conv.append({"role": "assistant", "content": assistant_content})
        tool_results: list[dict[str, Any]] = []
        for b in assistant_content:
            if b.get("type") != "tool_use":
                continue
            args = b.get("input") or {}
            if not isinstance(args, dict):
                args = {}
            result, events = _execute_tool(b["name"], args)
            for e in events:
                yield e
            tool_results.append({"type": "tool_result", "tool_use_id": b["id"], "content": result})
        conv.append({"role": "user", "content": tool_results})

    yield _nd({"error": "Reached the tool-call limit.", "done": True})


def _stream_openai_compatible(
    base: str, label: str, model: str, messages: list[dict[str, Any]], options: Any, key: str,
    tools: list[dict[str, Any]] | None = None,
) -> Iterator[bytes]:
    """Stream any OpenAI-compatible Chat Completions endpoint (OpenAI, Grok/xAI,
    Gemini's compat endpoint), translated into the app's NDJSON shape. When `tools`
    is given, runs the tool_calls loop: stream a turn, run any requested tools via
    _execute_tool, feed the tool messages back, and repeat until the model stops
    requesting tools (or MAX_TOOL_ROUNDS). The schemas are already OpenAI-shaped, so
    they pass through verbatim."""
    conv = [
        {"role": m["role"], "content": m["content"]}
        for m in messages
        if m.get("role") in ("system", "user", "assistant") and isinstance(m.get("content"), str) and m.get("content")
    ]
    in_tok = out_tok = 0
    for _round in range(MAX_TOOL_ROUNDS if tools else 1):
        body: dict[str, Any] = {
            "model": model, "messages": conv, "stream": True,
            "stream_options": {"include_usage": True},
        }
        # Temperature is intentionally omitted for ALL cloud models — newer ones
        # (GPT-5.x and others) reject any non-default value ("only the default (1) is
        # supported"). Let each provider use its own default.
        if tools:
            body["tools"] = tools
        calls: dict[int, dict[str, str]] = {}  # index -> {id, name, args}
        finish_reason = None
        round_in = round_out = 0
        try:
            with requests.post(
                base + "/chat/completions", json=body, stream=True, timeout=None,
                headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
            ) as up:
                if up.status_code >= 400:
                    detail = _cloud_err(up, label)
                    low = detail.lower()
                    # Some OpenAI models aren't served on chat/completions. Fall back to
                    # whichever endpoint the error names: the legacy Completions API, or
                    # the Responses API (the *-pro reasoning models). Route by the hint —
                    # don't send a v1/completions model to v1/responses.
                    if "api.openai.com" in base:
                        if "v1/completions" in low:
                            yield from _stream_openai_legacy_completions(base, label, model, conv, key)
                            return
                        if "v1/responses" in low or "not a chat model" in low:
                            yield from _stream_openai_responses(base, label, model, conv, key)
                            return
                    yield _nd({"error": detail, "done": True})
                    return
                for raw in up.iter_lines():
                    if not raw:
                        continue
                    line = raw.decode("utf-8", "replace").strip()
                    if not line.startswith("data:"):
                        continue
                    data = line[5:].strip()
                    if data == "[DONE]":
                        break
                    try:
                        ev = json.loads(data)
                    except ValueError:
                        continue
                    ch = ev.get("choices") or []
                    if ch:
                        c0 = ch[0]
                        delta = c0.get("delta", {}) or {}
                        if delta.get("content"):
                            yield _nd({"message": {"content": delta["content"]}})
                        for tc in (delta.get("tool_calls") or []):
                            slot = calls.setdefault(tc.get("index", 0), {"id": "", "name": "", "args": ""})
                            if tc.get("id"):
                                slot["id"] = tc["id"]
                            fn = tc.get("function") or {}
                            if fn.get("name"):
                                slot["name"] = fn["name"]
                            if fn.get("arguments"):
                                slot["args"] += fn["arguments"]
                        if c0.get("finish_reason"):
                            finish_reason = c0["finish_reason"]
                    usage = ev.get("usage")
                    if usage:
                        round_in = usage.get("prompt_tokens", round_in)
                        round_out = usage.get("completion_tokens", round_out)
        except requests.RequestException as exc:
            yield _nd({"error": f"{label} request failed: {exc}", "done": True})
            return

        in_tok += round_in
        out_tok += round_out
        named = {i: c for i, c in calls.items() if c.get("name")}
        if not tools or finish_reason != "tool_calls" or not named:
            yield _nd({"done": True, "prompt_eval_count": in_tok, "eval_count": out_tok})
            return

        # Re-create the assistant turn carrying the tool_calls, then run each tool and
        # feed its output back as a tool-role message keyed by tool_call_id.
        assistant_tcs: list[dict[str, Any]] = []
        for i in sorted(named):
            c = named[i]
            assistant_tcs.append({
                "id": c["id"] or f"call_{i}", "type": "function",
                "function": {"name": c["name"], "arguments": c["args"] or "{}"},
            })
        conv.append({"role": "assistant", "content": None, "tool_calls": assistant_tcs})
        for tc in assistant_tcs:
            try:
                args = json.loads(tc["function"]["arguments"] or "{}")
            except ValueError:
                args = {}
            if not isinstance(args, dict):
                args = {}
            result, events = _execute_tool(tc["function"]["name"], args)
            for e in events:
                yield e
            conv.append({"role": "tool", "tool_call_id": tc["id"], "content": result})

    yield _nd({"error": "Reached the tool-call limit.", "done": True})


def _stream_openai_legacy_completions(
    base: str, label: str, model: str, conv: list[dict[str, Any]], key: str,
) -> Iterator[bytes]:
    """Stream the legacy /v1/completions endpoint for completion-only models that
    reject chat/completions ('Did you mean to use v1/completions?'). The conversation
    is flattened into a single prompt. Text-only."""
    parts: list[str] = []
    for m in conv:
        role, content = m.get("role"), m.get("content", "")
        if not isinstance(content, str) or not content:
            continue
        if role == "system":
            parts.append(content)
        elif role == "user":
            parts.append("User: " + content)
        elif role == "assistant":
            parts.append("Assistant: " + content)
    prompt = "\n\n".join(parts)
    if conv and conv[-1].get("role") == "user":
        prompt += "\n\nAssistant:"
    body = {"model": model, "prompt": prompt, "stream": True, "max_tokens": 4096}
    in_tok = out_tok = 0
    try:
        with requests.post(
            base + "/completions", json=body, stream=True, timeout=None,
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
        ) as up:
            if up.status_code >= 400:
                yield _nd({"error": _cloud_err(up, label), "done": True})
                return
            for raw in up.iter_lines():
                if not raw:
                    continue
                line = raw.decode("utf-8", "replace").strip()
                if not line.startswith("data:"):
                    continue
                data = line[5:].strip()
                if data == "[DONE]":
                    break
                try:
                    ev = json.loads(data)
                except ValueError:
                    continue
                ch = ev.get("choices") or []
                if ch and ch[0].get("text"):
                    yield _nd({"message": {"content": ch[0]["text"]}})
                usage = ev.get("usage")
                if usage:
                    in_tok = usage.get("prompt_tokens", in_tok)
                    out_tok = usage.get("completion_tokens", out_tok)
            yield _nd({"done": True, "prompt_eval_count": in_tok, "eval_count": out_tok})
    except requests.RequestException as exc:
        yield _nd({"error": f"{label} request failed: {exc}", "done": True})


def _stream_openai_responses(
    base: str, label: str, model: str, conv: list[dict[str, Any]], key: str,
) -> Iterator[bytes]:
    """Stream OpenAI's Responses API (/v1/responses) for models that aren't served on
    chat/completions (the *-pro reasoning models). Text-only — tool-calling isn't
    wired on this path. Translated into the app's NDJSON shape."""
    system = ""
    inp: list[dict[str, Any]] = []
    for m in conv:
        role, content = m.get("role"), m.get("content", "")
        if not isinstance(content, str) or not content:
            continue
        if role == "system":
            system += ("\n\n" if system else "") + content
        elif role in ("user", "assistant"):
            inp.append({"role": role, "content": content})
    body: dict[str, Any] = {"model": model, "input": inp, "stream": True}
    if system:
        body["instructions"] = system
    in_tok = out_tok = 0
    try:
        with requests.post(
            base + "/responses", json=body, stream=True, timeout=None,
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"},
        ) as up:
            if up.status_code >= 400:
                yield _nd({"error": _cloud_err(up, label), "done": True})
                return
            for raw in up.iter_lines():
                if not raw:
                    continue
                line = raw.decode("utf-8", "replace").strip()
                if not line.startswith("data:"):
                    continue
                data = line[5:].strip()
                if data == "[DONE]":
                    break
                try:
                    ev = json.loads(data)
                except ValueError:
                    continue
                t = ev.get("type", "")
                if t == "response.output_text.delta":
                    if ev.get("delta"):
                        yield _nd({"message": {"content": ev["delta"]}})
                elif t in ("response.completed", "response.incomplete"):
                    u = (ev.get("response", {}) or {}).get("usage", {}) or {}
                    in_tok = u.get("input_tokens", in_tok)
                    out_tok = u.get("output_tokens", out_tok)
                elif t in ("response.failed", "error"):
                    err = (ev.get("response", {}) or {}).get("error") or ev.get("error") or {}
                    yield _nd({"error": (err.get("message") if isinstance(err, dict) else str(err)) or "Responses API error.", "done": True})
                    return
            yield _nd({"done": True, "prompt_eval_count": in_tok, "eval_count": out_tok})
    except requests.RequestException as exc:
        yield _nd({"error": f"{label} request failed: {exc}", "done": True})


# Real context-window sizes for cloud models (the UI's context meter max). Matched by
# id substring; falls back to a per-provider default. Approximate but realistic.
def _cloud_context_window(provider: str, model: str) -> int:
    m = (model or "").lower()
    if provider == "openai":
        if "gpt-4.1" in m:
            return 1_047_576
        if "gpt-5" in m or m.startswith("o1") or m.startswith("o3") or m.startswith("o4") or "-o1" in m or "-o3" in m:
            return 400_000 if "gpt-5" in m else 200_000
        if "gpt-4o" in m or "gpt-4-turbo" in m or m == "gpt-4" or "gpt-4-" in m:
            return 128_000
        if "gpt-3.5" in m:
            return 16_385
        return 128_000
    if provider == "anthropic":
        return 200_000  # Claude 3.x/4.x standard window
    if provider == "gemini":
        if "1.5-pro" in m:
            return 2_097_152
        if "flash" in m or "2.5" in m or "2.0" in m or "pro" in m:
            return 1_048_576
        return 1_048_576
    if provider == "grok":
        if "grok-4" in m:
            return 256_000
        return 131_072
    return 8_192


def _stream_cloud(
    provider: str, model: str, messages: list[dict[str, Any]], options: Any,
    tools: list[dict[str, Any]] | None = None,
) -> Iterator[bytes]:
    """Route a cloud chat to the right provider. The real send tool (key) stays
    server-side; the model id is whatever the user's account exposes. `tools`, when
    present, is the same read-only tool set the local models get — each provider's
    streamer runs its own tool loop over it."""
    key = PROVIDER_KEYS.get(provider, "")
    if not key:
        yield _nd({"error": f"{provider} is not connected. Add a key in Settings → Cloud models.", "done": True})
        return
    # Tell the UI this model's real context window so the meter's max is accurate.
    yield _nd({"ctx_window": _cloud_context_window(provider, model)})
    if provider == "anthropic":
        yield from _stream_anthropic(model, messages, options, key, tools)
    elif provider in OPENAI_COMPAT_BASE:
        labels = {"openai": "OpenAI", "grok": "Grok", "gemini": "Gemini"}
        yield from _stream_openai_compatible(
            OPENAI_COMPAT_BASE[provider], labels.get(provider, provider), model, messages, options, key, tools
        )
    else:
        yield _nd({"error": f"Unknown provider: {provider}", "done": True})


def _cloud_err(resp: requests.Response, label: str) -> str:
    """Pull a human message out of a failed provider response."""
    try:
        return resp.json().get("error", {}).get("message", "") or f"{label} returned {resp.status_code}."
    except ValueError:
        return (resp.text or "")[:300] or f"{label} returned {resp.status_code}."


def _gemini_inline_image(payload: dict[str, Any]) -> tuple[str, str] | None:
    """Find the first inline image (base64, mime) in a Gemini generateContent or
    predict response, tolerating both camelCase and snake_case shapes."""
    for cand in payload.get("candidates", []) or []:
        for part in ((cand.get("content") or {}).get("parts") or []):
            inline = part.get("inlineData") or part.get("inline_data")
            if inline and inline.get("data"):
                return inline["data"], (inline.get("mimeType") or inline.get("mime_type") or "image/png")
    for pred in payload.get("predictions", []) or []:
        if pred.get("bytesBase64Encoded"):
            return pred["bytesBase64Encoded"], pred.get("mimeType", "image/png")
    return None


def _gemini_image(model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Gemini image generation. Imagen models use :predict; native gemini-*-image
    models use :generateContent with IMAGE response modality. Renders inline."""
    yield _nd({"message": {"content": "🎨 Generating image…\n\n"}})
    try:
        if "imagen" in model.lower():
            r = requests.post(
                f"{GEMINI_NATIVE_BASE}/models/{model}:predict", params={"key": key},
                json={"instances": [{"prompt": prompt}], "parameters": {"sampleCount": 1}}, timeout=180,
            )
        else:
            r = requests.post(
                f"{GEMINI_NATIVE_BASE}/models/{model}:generateContent", params={"key": key},
                json={"contents": [{"parts": [{"text": prompt}]}],
                      "generationConfig": {"responseModalities": ["TEXT", "IMAGE"]}}, timeout=180,
            )
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Gemini image"), "done": True})
            return
        found = _gemini_inline_image(r.json())
        if not found:
            yield _nd({"error": "Gemini returned no image.", "done": True})
            return
        b64, mime = found
        url = _save_media(base64.b64decode(b64), ".png" if "png" in mime else ".jpg")
        yield _nd({"media": {"kind": "image", "url": url, "alt": prompt[:120]}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Gemini image request failed: {exc}", "done": True})


def _gemini_video(model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Gemini (Veo) video generation — a long-running operation: submit, poll the
    operation until done, then download the returned video URI. Renders inline."""
    yield _nd({"message": {"content": "🎬 Submitting Veo video job — this can take a minute…\n\n"}})
    try:
        r = requests.post(
            f"{GEMINI_NATIVE_BASE}/models/{model}:predictLongRunning", params={"key": key},
            json={"instances": [{"prompt": prompt}], "parameters": {"sampleCount": 1}}, timeout=60,
        )
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Gemini video"), "done": True})
            return
        op = r.json().get("name")
        if not op:
            yield _nd({"error": "Veo job did not return an operation name.", "done": True})
            return
        waited, done, data = 0, False, {}
        while not done and waited < VIDEO_MAX_WAIT_SECS:
            time.sleep(VIDEO_POLL_SECS)
            waited += VIDEO_POLL_SECS
            try:
                pr = requests.get(f"{GEMINI_NATIVE_BASE}/{op}", params={"key": key}, timeout=30)
                if pr.status_code < 400:
                    data = pr.json()
                    done = bool(data.get("done"))
            except requests.RequestException:
                pass
            yield _nd({"message": {"content": f"…{waited}s — {'done' if done else 'rendering'}\n"}})
        if not done:
            yield _nd({"error": f"Video generation didn't finish in {VIDEO_MAX_WAIT_SECS}s.", "done": True})
            return
        if data.get("error"):
            err = data["error"]
            yield _nd({"error": str(err.get("message", err) if isinstance(err, dict) else err), "done": True})
            return
        resp = data.get("response", {}) or {}
        samples = (((resp.get("generateVideoResponse") or {}).get("generatedSamples"))
                   or resp.get("generatedSamples") or [])
        uri = None
        for s in samples:
            vid = s.get("video") or {}
            uri = vid.get("uri") or vid.get("url")
            if uri:
                break
        if not uri:
            yield _nd({"error": "Veo returned no video URI.", "done": True})
            return
        dl = requests.get(uri, params=None if "key=" in uri else {"key": key}, timeout=180)
        if dl.status_code >= 400:
            yield _nd({"error": _cloud_err(dl, "Veo download"), "done": True})
            return
        url = _save_media(dl.content, ".mp4")
        yield _nd({"media": {"kind": "video", "url": url, "mime": "video/mp4"}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Gemini video request failed: {exc}", "done": True})


def _generate_image(provider: str, model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Image generation. Gemini uses its native API; OpenAI/xAI use the
    OpenAI-compatible /images/generations endpoint. Renders inline as an image."""
    if provider == "gemini":
        yield from _gemini_image(model, prompt, key)
        return
    base = OPENAI_COMPAT_BASE.get(provider)
    if not base:
        yield _nd({"error": f"Image generation isn't supported for {provider} in this app yet.", "done": True})
        return
    yield _nd({"message": {"content": "🎨 Generating image…\n\n"}})
    try:
        r = requests.post(
            base + "/images/generations", json={"model": model, "prompt": prompt, "n": 1},
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"}, timeout=180,
        )
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Image"), "done": True})
            return
        data = (r.json().get("data") or [])
        if not data:
            yield _nd({"error": "No image was returned.", "done": True})
            return
        item = data[0]
        if item.get("b64_json"):
            url = _save_media(base64.b64decode(item["b64_json"]), ".png")
        elif item.get("url"):
            url = item["url"]  # remote provider URL; rendered directly
        else:
            yield _nd({"error": "No image data in the response.", "done": True})
            return
        yield _nd({"media": {"kind": "image", "url": url, "alt": prompt[:120]}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Image request failed: {exc}", "done": True})


def _wav_wrap(pcm: bytes, rate: int = 24000, channels: int = 1, sample_width: int = 2) -> bytes:
    """Wrap raw little-endian PCM (Gemini TTS returns L16) in a minimal WAV header
    so browsers can play it (they can't play bare PCM)."""
    import struct
    byte_rate = rate * channels * sample_width
    block_align = channels * sample_width
    return (
        b"RIFF" + struct.pack("<I", 36 + len(pcm)) + b"WAVE"
        + b"fmt " + struct.pack("<IHHIIHH", 16, 1, channels, rate, byte_rate, block_align, sample_width * 8)
        + b"data" + struct.pack("<I", len(pcm)) + pcm
    )


def _gemini_speech(model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Gemini TTS via :generateContent with AUDIO modality. Returns raw PCM (L16),
    which we wrap as WAV and render inline."""
    yield _nd({"message": {"content": "🔊 Synthesizing speech…\n\n"}})
    try:
        r = requests.post(
            f"{GEMINI_NATIVE_BASE}/models/{model}:generateContent", params={"key": key},
            json={"contents": [{"parts": [{"text": prompt}]}],
                  "generationConfig": {"responseModalities": ["AUDIO"],
                      "speechConfig": {"voiceConfig": {"prebuiltVoiceConfig": {"voiceName": "Kore"}}}}},
            timeout=120,
        )
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Gemini speech"), "done": True})
            return
        b64, mime = None, ""
        for cand in r.json().get("candidates", []) or []:
            for part in ((cand.get("content") or {}).get("parts") or []):
                inline = part.get("inlineData") or part.get("inline_data")
                if inline and inline.get("data"):
                    b64 = inline["data"]
                    mime = inline.get("mimeType") or inline.get("mime_type") or ""
                    break
            if b64:
                break
        if not b64:
            yield _nd({"error": "Gemini returned no audio.", "done": True})
            return
        rate_m = re.search(r"rate=(\d+)", mime or "")
        wav = _wav_wrap(base64.b64decode(b64), rate=int(rate_m.group(1)) if rate_m else 24000)
        url = _save_media(wav, ".wav")
        yield _nd({"media": {"kind": "audio", "url": url, "mime": "audio/wav"}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Gemini speech request failed: {exc}", "done": True})


def _generate_speech(provider: str, model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Text-to-speech. Gemini uses its native AUDIO-modality API; OpenAI uses
    /audio/speech. Speaks the prompt and renders an inline audio player."""
    if provider == "gemini":
        yield from _gemini_speech(model, prompt, key)
        return
    base = OPENAI_COMPAT_BASE.get(provider)
    if provider != "openai" or not base:
        yield _nd({"error": f"Text-to-speech isn't supported for {provider} in this app yet.", "done": True})
        return
    yield _nd({"message": {"content": "🔊 Synthesizing speech…\n\n"}})
    try:
        r = requests.post(
            base + "/audio/speech",
            json={"model": model, "input": prompt, "voice": "alloy", "response_format": "mp3"},
            headers={"Authorization": f"Bearer {key}", "Content-Type": "application/json"}, timeout=120,
        )
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Speech"), "done": True})
            return
        url = _save_media(r.content, ".mp3")
        yield _nd({"media": {"kind": "audio", "url": url, "mime": "audio/mpeg"}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Speech request failed: {exc}", "done": True})


def _generate_video(provider: str, model: str, prompt: str, key: str) -> Iterator[bytes]:
    """Video generation. Gemini uses Veo's long-running operation API; OpenAI/xAI
    use the async /videos jobs API (Sora pattern; xAI best-effort). Submit → poll →
    download, streaming progress, then render inline."""
    if provider == "gemini":
        yield from _gemini_video(model, prompt, key)
        return
    base = OPENAI_COMPAT_BASE.get(provider)
    if not base:
        yield _nd({"error": f"Video generation isn't supported for {provider} in this app yet.", "done": True})
        return
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    yield _nd({"message": {"content": "🎬 Submitting video job — this can take a minute…\n\n"}})
    try:
        r = requests.post(base + "/videos", json={"model": model, "prompt": prompt}, headers=headers, timeout=60)
        if r.status_code >= 400:
            yield _nd({"error": _cloud_err(r, "Video"), "done": True})
            return
        job = r.json()
        vid = job.get("id")
        status = (job.get("status") or "").lower()
        if not vid:
            yield _nd({"error": "Video job did not return an id.", "done": True})
            return
        waited = 0
        done_states = {"completed", "succeeded", "success"}
        fail_states = {"failed", "error", "cancelled", "canceled"}
        while status not in done_states and status not in fail_states and waited < VIDEO_MAX_WAIT_SECS:
            time.sleep(VIDEO_POLL_SECS)
            waited += VIDEO_POLL_SECS
            try:
                pr = requests.get(base + f"/videos/{vid}", headers=headers, timeout=30)
                if pr.status_code < 400:
                    job = pr.json()
                    status = (job.get("status") or status).lower()
            except requests.RequestException:
                pass
            prog = job.get("progress")
            tail = f" ({prog}%)" if isinstance(prog, (int, float)) else ""
            yield _nd({"message": {"content": f"…{waited}s — {status or 'working'}{tail}\n"}})
        if status not in done_states:
            yield _nd({"error": f"Video generation didn't finish (status: {status or 'timed out'}).", "done": True})
            return
        cr = requests.get(base + f"/videos/{vid}/content", headers=headers, timeout=180)
        if cr.status_code >= 400:
            yield _nd({"error": _cloud_err(cr, "Video download"), "done": True})
            return
        url = _save_media(cr.content, ".mp4")
        yield _nd({"media": {"kind": "video", "url": url, "mime": "video/mp4"}})
        yield _nd({"done": True})
    except requests.RequestException as exc:
        yield _nd({"error": f"Video request failed: {exc}", "done": True})


def _stream_cloud_media(provider: str, model: str, modality: str, prompt: str) -> Iterator[bytes]:
    """Drive a non-chat cloud model (image / speech / video) from the user's prompt,
    or explain when a model type can't be driven as chat input."""
    key = PROVIDER_KEYS.get(provider, "")
    if not key:
        yield _nd({"error": f"{provider} is not connected. Add a key in Settings → Cloud models.", "done": True})
        return
    if not (prompt or "").strip():
        yield _nd({"error": "Type a prompt describing what to generate.", "done": True})
        return
    if modality == "image":
        yield from _generate_image(provider, model, prompt, key)
    elif modality == "audio":
        yield from _generate_speech(provider, model, prompt, key)
    elif modality == "video":
        yield from _generate_video(provider, model, prompt, key)
    else:  # "unsupported" — transcription, realtime, embeddings, moderation, …
        yield _nd({"message": {"content": (
            f"⚠️ **{model}** is a `{modality}`-type model — it takes audio/other input "
            "rather than a chat prompt, so there's nothing to run here. Pick a chat, "
            "image, voice, or video model instead."
        )}})
        yield _nd({"done": True})


def _stream_chat(body: dict[str, Any]) -> Iterator[bytes]:
    """Proxy a streaming chat request to Ollama, yielding NDJSON lines verbatim.

    On any transport failure, emit a single NDJSON error object so the browser
    can surface it inside the transcript instead of failing silently.
    """
    try:
        with requests.post(
            f"{OLLAMA_HOST}/api/chat",
            json=body,
            stream=True,
            timeout=None,  # streaming generations can run long
        ) as upstream:
            if upstream.status_code >= 400:
                # Surface Ollama's own message — e.g. asking a non-reasoning model
                # to think returns 400 "<model> does not support thinking".
                try:
                    detail = upstream.json().get("error", "")
                except ValueError:
                    detail = (upstream.text or "").strip()[:300]
                error = {"error": detail or f"Ollama returned {upstream.status_code}.", "done": True}
                yield (json.dumps(error) + "\n").encode("utf-8")
                return
            for line in upstream.iter_lines():
                if line:
                    yield line + b"\n"
    except requests.RequestException as exc:
        error = {"error": f"Stream failed: {exc}", "done": True}
        yield (json.dumps(error) + "\n").encode("utf-8")


def _chat_with_tools(
    model: str, messages: list[dict[str, Any]], options: Any, think: bool,
    tools: list[dict[str, Any]],
) -> Iterator[bytes]:
    """Stream a chat that may call read-only MCP tools. Each round streams content/
    thinking verbatim; if the model requested tools, run them (read-only, so nothing
    is mutated), append the results, and loop. `{"tool_event": …}` lines let the UI
    show which tools ran. Falls back cleanly on any transport error."""
    convo = list(messages)
    for _round in range(MAX_TOOL_ROUNDS):
        body: dict[str, Any] = {
            "model": model, "messages": convo, "stream": True, "tools": tools,
            "keep_alive": KEEP_ALIVE,
        }
        if options:
            body["options"] = options
        if think:
            body["think"] = True

        tool_calls: list[dict[str, Any]] = []
        assistant_text = ""
        try:
            with requests.post(
                f"{OLLAMA_HOST}/api/chat", json=body, stream=True, timeout=None
            ) as up:
                if up.status_code >= 400:
                    try:
                        detail = up.json().get("error", "")
                    except ValueError:
                        detail = (up.text or "").strip()[:300]
                    yield (json.dumps({"error": detail or f"Ollama returned {up.status_code}.", "done": True}) + "\n").encode("utf-8")
                    return
                for raw in up.iter_lines():
                    if not raw:
                        continue
                    yield raw + b"\n"  # pass content + thinking to the client verbatim
                    try:
                        obj = json.loads(raw)
                    except ValueError:
                        continue
                    msg = obj.get("message") or {}
                    if msg.get("content"):
                        assistant_text += msg["content"]
                    if msg.get("tool_calls"):
                        tool_calls.extend(msg["tool_calls"])
        except requests.RequestException as exc:
            yield (json.dumps({"error": f"Stream failed: {exc}", "done": True}) + "\n").encode("utf-8")
            return

        if not tool_calls:
            return  # no tool requested — the final answer already streamed

        convo.append({"role": "assistant", "content": assistant_text, "tool_calls": tool_calls})
        for tc in tool_calls:
            fn = tc.get("function") or {}
            name = fn.get("name", "")
            args = fn.get("arguments", {})
            if isinstance(args, str):
                try:
                    args = json.loads(args)
                except ValueError:
                    args = {}
            if not isinstance(args, dict):
                args = {}
            result, events = _execute_tool(name, args)
            for e in events:
                yield e
            convo.append({"role": "tool", "content": result, "tool_name": name})

    yield _nd({"error": "Reached the tool-call limit.", "done": True})


@app.post("/api/chat")
def chat() -> Response:
    """Accept {model, messages, options?} and stream Ollama's reply back as NDJSON."""
    data: dict[str, Any] | None = request.get_json(silent=True)
    if not data or "model" not in data or "messages" not in data:
        return jsonify({"error": "Body must include 'model' and 'messages'."}), 400

    messages = list(data["messages"])

    # Ground the answer in the user's second brain using their latest message.
    last_user = next(
        (m.get("content", "") for m in reversed(messages)
         if isinstance(m, dict) and m.get("role") == "user"),
        "",
    )
    # Cloud models get the same grounding, MINUS private/ PII — it would otherwise be sent
    # to the provider. Local/offline models retrieve private notes normally.
    _is_cloud_req = str(data.get("provider", "")).lower() in PROVIDER_CONFIG_KEY
    context, grounding = "", []
    if isinstance(last_user, str) and last_user.strip():
        context, grounding = _retrieve_context(last_user, exclude_private=_is_cloud_req)

    # Native tools (Brave/weather) + read-only MCP tools (Monarch/Gmail/files).
    tools = _native_tool_schemas() + MCP.ollama_tools()
    # A draft-only iMessage tool (sends nothing; the user confirms in the UI).
    if MCP.imessage_available():
        tools = tools + [DRAFT_IMESSAGE_SCHEMA]

    # Fold the current date (the model's training cutoff makes it guess stale
    # years, which breaks date-filtered tool calls) and the retrieved context into
    # the system message. The base prompt (instructions.md) already covers tool
    # usage. All stays server-side, so none of it bloats saved chats.
    date_line = (
        "\n\nThe current date and time is "
        + datetime.now().astimezone().strftime("%A, %B %d, %Y at %I:%M %p %Z") + "."
    )
    system_extra = date_line + context
    # Inject the open shared document so every model can READ it (only the
    # pen-holder gets the edit tool, below).
    if OPEN_DOC["path"]:
        system_extra += (
            "\n\n# Shared document (live viewer)\nThe user has this document open. You "
            "can read and reference it. Only edit it if you have the edit_document tool.\n\n"
            "--- " + os.path.basename(OPEN_DOC["path"]) + " ---\n" + OPEN_DOC["content"][:8000]
        )
    if messages and isinstance(messages[0], dict) and messages[0].get("role") == "system":
        messages[0] = {**messages[0], "content": messages[0].get("content", "") + system_extra}
    else:
        messages = [{"role": "system", "content": SYSTEM_PROMPT + system_extra}, *messages]

    model = data["model"]
    # Cloud routing: a connected provider ("anthropic"/"openai") sends this turn to
    # that API instead of Ollama. RAG grounding above is still applied (the user
    # opted to ground cloud models too). Cloud turns skip Ollama tools/num_ctx.
    provider = str(data.get("provider", "")).lower()
    # Any known cloud provider routes to the cloud path — even with no key yet, so a
    # pre-listed flagship returns a clean "connect your key" message (not an Ollama 404).
    # Cloud turns get the SAME read-only tool set as local (Brave/weather/MCP/draft),
    # run through each provider's own tool loop. They skip num_ctx (no Ollama window).
    is_cloud = provider in PROVIDER_CONFIG_KEY
    # Set num_ctx explicitly (capped) so the window is deterministic and the UI's
    # context meter is truthful — Ollama's default would otherwise truncate silently.
    options = {**(data.get("options") or {})}
    if not is_cloud:
        options["num_ctx"] = _effective_num_ctx(model)
    think = bool(data.get("think"))

    # The pen: only the selected pane's request carries can_edit_doc, so only that
    # model is handed the edit tool. Cloud models hold the pen too — their tool loops
    # run edit_document the same way the local loop does.
    if data.get("can_edit_doc") and OPEN_DOC["path"] and OPEN_DOC.get("editable"):
        tools = tools + [EDIT_DOCUMENT_SCHEMA]

    body: dict[str, Any] = {
        "model": model, "messages": messages, "stream": True, "keep_alive": KEEP_ALIVE,
    }
    if options:
        body["options"] = options
    # Reasoning models (deepseek-r1, qwen3, …) stream a separate message.thinking
    # field; the UI sets this automatically from /api/capabilities.
    if think:
        body["think"] = True

    def generate() -> Iterator[bytes]:
        # Non-chat cloud models (image / speech / video) generate from the raw
        # prompt — skip RAG grounding and the chat streamer entirely.
        if is_cloud:
            modality = _cloud_modality(model)
            if modality != "chat":
                yield from _stream_cloud_media(provider, model, modality, last_user if isinstance(last_user, str) else "")
                return
        # A leading metadata line lets the UI show what notes grounded the answer.
        if grounding:
            yield (json.dumps({"grounding": grounding}) + "\n").encode("utf-8")
            yield (json.dumps({"rag_meta": {
                "k": RAG_TOP_K, "returned": len(grounding), "min_score": RAG_MIN_SCORE,
            }}) + "\n").encode("utf-8")
        if is_cloud:
            yield from _stream_cloud(provider, model, messages, options, tools)
        elif tools and _model_supports_tools(model):
            yield from _chat_with_tools(model, messages, options, think, tools)
        else:
            yield from _stream_chat(body)

    return Response(generate(), mimetype="application/x-ndjson")


@app.get("/api/providers")
def api_providers() -> Response:
    """Cloud-provider connection status + each connected provider's model list.
    Never returns the key itself. Lazily fills the model cache for a provider
    that has a key but no cached list yet (e.g. a key loaded from config at boot)."""
    out: dict[str, Any] = {}
    for p in PROVIDER_CONFIG_KEY:
        key = PROVIDER_KEYS.get(p, "")
        if key and not PROVIDER_MODELS.get(p):
            try:
                PROVIDER_MODELS[p] = _list_provider_models(p, key)
            except Exception:
                pass
        full = PROVIDER_MODELS.get(p, [])
        curated = CURATED_CLOUD.get(p, [])
        # Only the curated top-3 the account actually exposes (in curated order);
        # if connected but none match (e.g. a renamed model), fall back to its first 3.
        if full:
            shown = [m for m in curated if m in set(full)] or full[:3]
        else:
            shown = curated  # placeholders until the key is connected
        out[p] = {"connected": bool(key), "models": shown}
    return jsonify(out)


@app.post("/api/providers/connect")
def api_provider_connect() -> Response:
    """Validate a pasted API key by listing the account's models; on success,
    store it (gitignored config.local.json) and cache the model list. The key is
    never echoed back."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    provider = str(data.get("provider", "")).lower()
    key = str(data.get("key", "")).strip()
    if provider not in PROVIDER_CONFIG_KEY:
        return jsonify({"ok": False, "error": "Unknown provider."}), 400
    if not key:
        return jsonify({"ok": False, "error": "API key is required."}), 400
    try:
        models = _list_provider_models(provider, key)
    except Exception as exc:
        return jsonify({"ok": False, "error": f"Could not validate key: {exc}"})
    PROVIDER_KEYS[provider] = key
    PROVIDER_MODELS[provider] = models
    _save_local_config_key(PROVIDER_CONFIG_KEY[provider], key)
    return jsonify({"ok": True, "connected": True, "models": models})


@app.post("/api/providers/disconnect")
def api_provider_disconnect() -> Response:
    """Forget a provider's key (in memory + config.local.json)."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    provider = str(data.get("provider", "")).lower()
    if provider not in PROVIDER_CONFIG_KEY:
        return jsonify({"ok": False, "error": "Unknown provider."}), 400
    PROVIDER_KEYS[provider] = ""
    PROVIDER_MODELS[provider] = []
    _save_local_config_key(PROVIDER_CONFIG_KEY[provider], "")
    return jsonify({"ok": True, "connected": False})


@app.get("/api/doc")
def api_doc() -> Response:
    """Current open document (the live viewer reads this)."""
    p = OPEN_DOC["path"]
    return jsonify({
        "open": bool(p), "path": p,
        "name": os.path.basename(p) if p else "",
        "content": OPEN_DOC["content"],
        "kind": OPEN_DOC.get("kind", ""),
        "editable": OPEN_DOC.get("editable", False),
    })


@app.post("/api/doc/open")
def api_doc_open() -> Response:
    """Open a file (any type) by path into the shared live viewer."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    return _open_doc_path(str(data.get("path", "")))


@app.post("/api/doc/browse")
def api_doc_browse() -> Response:
    """Pop the native macOS file picker so the user can choose ANY file on their
    computer, then open it. Blocks until they pick or cancel."""
    try:
        proc = subprocess.run(
            ["osascript", "-e", 'POSIX path of (choose file with prompt "Open a file in Local LLM Studio")'],
            capture_output=True, text=True, timeout=300,
        )
    except (subprocess.SubprocessError, OSError) as exc:
        return jsonify({"ok": False, "error": f"File picker unavailable: {exc}"})
    if proc.returncode != 0:  # user pressed Cancel (osascript exits non-zero)
        return jsonify({"ok": False, "cancelled": True})
    return _open_doc_path(proc.stdout.strip())


@app.get("/api/doc/raw")
def api_doc_raw() -> Response:
    """Serve the open document's bytes (for previewing images/PDF/video inline)."""
    p = OPEN_DOC["path"]
    if not p or not os.path.isfile(p):
        return jsonify({"error": "No document open."}), 404
    return send_file(p)


@app.post("/api/doc/close")
def api_doc_close() -> Response:
    OPEN_DOC.update({"path": "", "content": "", "kind": "", "editable": False})
    return jsonify({"ok": True})


@app.post("/api/doc/save")
def api_doc_save() -> Response:
    """Manual save — used when the user holds the pen and edits in the viewer."""
    data: dict[str, Any] = request.get_json(silent=True) or {}
    if not OPEN_DOC["path"]:
        return jsonify({"ok": False, "error": "No document is open."}), 400
    content = str(data.get("content", ""))
    try:
        with open(OPEN_DOC["path"], "w", encoding="utf-8") as f:
            f.write(content)
    except OSError as exc:
        return jsonify({"ok": False, "error": f"Could not write: {exc}"}), 400
    OPEN_DOC["content"] = content
    return jsonify({"ok": True})


@app.post("/api/imessage/send")
def imessage_send() -> Response:
    """Send an iMessage — the ONLY send path, reached only when the user taps Send
    on the confirmation card. The model can never call this; it can only draft."""
    data: dict[str, Any] | None = request.get_json(silent=True)
    to = str((data or {}).get("to") or "").strip()
    text = str((data or {}).get("text") or "").strip()
    if not to or not text:
        return jsonify({"error": "Both 'to' and 'text' are required."}), 400
    result = MCP.send_imessage(to, text)
    ok = not result.lower().startswith(("tool error", "tool call failed", "refused", "unknown"))
    return jsonify({"ok": ok, "result": result})


# ─────────────────────────────────────────────────────────────────────────
# Project filesystem API (Phase 6, slice 1) — a path-jailed, READ-ONLY view of
# the folder the agent works in. Foundation for the file-tree sidebar (#1),
# @file mentions (#8), and the (approval-gated) agentic file/exec tools (#2/#3).
# Every path resolves through a realpath jail that refuses to escape the root.
# ─────────────────────────────────────────────────────────────────────────
_PROJECT_LOCK = threading.Lock()
_DEFAULT_PROJECT_ROOT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "workspace")
PROJECT_ROOT: str = os.path.realpath(os.path.expanduser(
    os.environ.get("LLS_PROJECT_ROOT") or _LOCAL_CONFIG.get("LLS_PROJECT_ROOT") or _DEFAULT_PROJECT_ROOT
))
os.makedirs(PROJECT_ROOT, exist_ok=True)

_FS_SKIP_DIRS = {".git", "node_modules", "__pycache__", ".venv", "venv", ".mypy_cache",
                 ".pytest_cache", "graphify-out", "dist", "build", ".lls-agent"}
_FS_TREE_CAP = 4000        # max entries from /api/fs/tree
_FS_READ_CAP = 200_000     # max bytes from /api/fs/read
_FS_SEARCH_CAP = 50        # max hits from /api/fs/search


def _project_path(rel: str) -> "str | None":
    """Resolve `rel` inside PROJECT_ROOT; returns the absolute realpath, or None
    if it escapes the jail (callers 400 instead of touching outside files)."""
    root = os.path.realpath(PROJECT_ROOT)
    p = os.path.realpath(os.path.join(root, rel or ""))
    return p if (p == root or p.startswith(root + os.sep)) else None


@app.get("/api/project")
def api_project_get() -> Response:
    """Report the current project root."""
    root = PROJECT_ROOT
    return jsonify({"root": root, "name": os.path.basename(root.rstrip(os.sep)) or root,
                    "exists": os.path.isdir(root)})


@app.post("/api/project")
def api_project_set() -> Response:
    """Set the project root to any existing directory inside the user's home."""
    global PROJECT_ROOT
    data = request.get_json(silent=True) or {}
    raw = str(data.get("path", "")).strip()
    if not raw:
        return jsonify({"error": "Body must include 'path'."}), 400
    cand = os.path.realpath(os.path.expanduser(raw))
    home = os.path.realpath(os.path.expanduser("~"))
    if not (cand == home or cand.startswith(home + os.sep)):
        return jsonify({"error": "Project root must be inside your home folder."}), 400
    if not os.path.isdir(cand):
        return jsonify({"error": "Not a directory."}), 400
    with _PROJECT_LOCK:
        PROJECT_ROOT = cand
    return jsonify({"root": cand, "name": os.path.basename(cand.rstrip(os.sep)) or cand})


@app.get("/api/fs/tree")
def api_fs_tree() -> Response:
    """Flat, path-jailed listing of the project (rel path + type + size), capped,
    skipping noise dirs. The UI builds the tree view from this."""
    root = os.path.realpath(PROJECT_ROOT)
    out: list[dict[str, Any]] = []
    truncated = False
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = sorted(d for d in dirnames if d not in _FS_SKIP_DIRS and not d.startswith("."))
        for name in list(dirnames):
            rel = os.path.relpath(os.path.join(dirpath, name), root)
            out.append({"path": rel, "type": "dir"})
            if len(out) >= _FS_TREE_CAP:
                truncated = True
                break
        if truncated:
            break
        for name in sorted(filenames):
            if name == ".DS_Store":
                continue
            full = os.path.join(dirpath, name)
            rel = os.path.relpath(full, root)
            try:
                size = os.path.getsize(full)
            except OSError:
                size = 0
            out.append({"path": rel, "type": "file", "size": size})
            if len(out) >= _FS_TREE_CAP:
                truncated = True
                break
        if truncated:
            break
    return jsonify({"root": root, "entries": out, "truncated": truncated})


@app.get("/api/fs/read")
def api_fs_read() -> Response:
    """Return a project file's text content, path-jailed and size-capped."""
    rel = request.args.get("path", "")
    p = _project_path(rel)
    if not p:
        return jsonify({"error": "Path escapes the project root."}), 400
    if not os.path.isfile(p):
        return jsonify({"error": "Not a file."}), 404
    try:
        with open(p, "rb") as fh:
            raw = fh.read(_FS_READ_CAP + 1)
    except OSError as exc:
        return jsonify({"error": f"Could not read: {exc}"}), 500
    truncated = len(raw) > _FS_READ_CAP
    raw = raw[:_FS_READ_CAP]
    try:
        text = raw.decode("utf-8")
        binary = False
    except UnicodeDecodeError:
        text, binary = "", True
    return jsonify({"path": rel, "content": text, "binary": binary, "truncated": truncated})


@app.get("/api/fs/search")
def api_fs_search() -> Response:
    """Filename search under the project root for @file autocomplete — query is a
    case-insensitive substring of the relative path; capped."""
    q = request.args.get("q", "").strip().lower()
    root = os.path.realpath(PROJECT_ROOT)
    hits: list[str] = []
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in _FS_SKIP_DIRS and not d.startswith(".")]
        for name in filenames:
            if name == ".DS_Store":
                continue
            rel = os.path.relpath(os.path.join(dirpath, name), root)
            if not q or q in rel.lower():
                hits.append(rel)
                if len(hits) >= _FS_SEARCH_CAP:
                    return jsonify({"hits": sorted(hits)[:_FS_SEARCH_CAP], "truncated": True})
    return jsonify({"hits": sorted(hits), "truncated": False})


# ─────────────────────────────────────────────────────────────────────────
# Agentic file CRUD + change journal (Phase 6, slice 3) — lets a model create,
# update, move, and delete files/folders in the project. SAFE BY DESIGN:
#   • OFF by default (AGENT_WRITES_ENABLED) — the tools aren't even shown to the
#     model until the user turns on "Agent edits".
#   • Path-jailed to the project root (reuses _project_path).
#   • EVERY change is journaled and reversible: updates keep a checkpoint copy,
#     deletes move to a project trash — nothing is hard-erased, so any change the
#     model makes can be undone.
# One core serves both the model tools (every provider, via _dispatch_tool) and
# the /api/fs/* HTTP endpoints the UI calls, so behavior is identical everywhere.
# ─────────────────────────────────────────────────────────────────────────
import shutil
import uuid

AGENT_WRITES_ENABLED: bool = False
AGENT_APPROVAL_REQUIRED: bool = True  # when True, the model's edits are staged for Approve/Reject
_AGENT_LOCK = threading.Lock()
_CRUD_TOOL_NAMES = {"project_write", "project_create_folder", "project_move", "project_delete"}
_PENDING_CHANGES: dict[str, dict[str, Any]] = {}  # id -> {name, args}


def _change_summary(name: str, args: dict[str, Any]) -> str:
    if name == "project_write":
        return "Write " + str(args.get("path", "?"))
    if name == "project_create_folder":
        return "Create folder " + str(args.get("path", "?"))
    if name == "project_move":
        return "Move " + str(args.get("from", "?")) + " → " + str(args.get("to", "?"))
    if name == "project_delete":
        return "Delete " + str(args.get("path", "?"))
    return name


def _proposed_diff(args: dict[str, Any]) -> str:
    """Unified diff between a file's current content and the model's proposed content."""
    import difflib
    rel = str(args.get("path") or "")
    p = _project_path(rel)
    before = ""
    if p and os.path.isfile(p):
        try:
            with open(p, "rb") as fh:
                before = fh.read(_FS_READ_CAP).decode("utf-8")
        except (OSError, UnicodeDecodeError):
            before = ""
    after = str(args.get("content") or "")
    return "".join(difflib.unified_diff(
        before.splitlines(keepends=True), after.splitlines(keepends=True),
        fromfile=rel + " (current)", tofile=rel + " (proposed)", n=3))


def _stage_change(name: str, args: dict[str, Any]) -> "tuple[str, list[bytes]]":
    """Hold a model-proposed change and emit an Approve/Reject card. Applied only
    when the user approves (api_change_approve), never by the model."""
    cid = "ch" + uuid.uuid4().hex[:10]
    _PENDING_CHANGES[cid] = {"name": name, "args": args}
    card: dict[str, Any] = {"id": cid, "op": name, "summary": _change_summary(name, args)}
    if name == "project_write":
        card["diff"] = _proposed_diff(args)
    result = ("Staged this change for the user's approval — it is applied ONLY if they tap "
              "Approve; you cannot apply it yourself. Wait for their decision and don't repeat it.")
    return result, [_nd({"confirm_change": card})]


def _agent_dir() -> str:
    d = os.path.join(os.path.realpath(PROJECT_ROOT), ".lls-agent")
    os.makedirs(os.path.join(d, "checkpoints"), exist_ok=True)
    os.makedirs(os.path.join(d, "trash"), exist_ok=True)
    return d


def _journal_load() -> list[dict[str, Any]]:
    try:
        with open(os.path.join(_agent_dir(), "journal.json"), encoding="utf-8") as fh:
            data = json.load(fh)
            return data if isinstance(data, list) else []
    except (OSError, ValueError):
        return []


def _journal_save(entries: list[dict[str, Any]]) -> None:
    try:
        with open(os.path.join(_agent_dir(), "journal.json"), "w", encoding="utf-8") as fh:
            json.dump(entries, fh, indent=2)
    except OSError:
        pass


def _journal_add(entry: dict[str, Any]) -> None:
    with _AGENT_LOCK:
        entries = _journal_load()
        seq = 0
        for e in entries:
            try:
                seq = max(seq, int(str(e.get("id", "c0"))[1:]))
            except ValueError:
                pass
        entries.append({"id": f"c{seq + 1}", "undone": False,
                        "ts": datetime.now().astimezone().isoformat(timespec="seconds"), **entry})
        _journal_save(entries)


def _checkpoint_copy(p: str) -> str:
    name = uuid.uuid4().hex[:12] + "__" + os.path.basename(p)
    shutil.copy2(p, os.path.join(_agent_dir(), "checkpoints", name))
    return name


def _writes_guard() -> "str | None":
    if not AGENT_WRITES_ENABLED:
        return ("File edits are disabled. Ask the user to turn on 'Agent edits' in the "
                "Files panel before trying to create, change, move, or delete files.")
    return None


def _crud_write(args: dict[str, Any]) -> str:
    guard = _writes_guard()
    if guard:
        return guard
    rel = str(args.get("path") or "").strip()
    content = args.get("content")
    content = "" if content is None else str(content)
    p = _project_path(rel)
    if not p or p == os.path.realpath(PROJECT_ROOT):
        return f"Refused: '{rel}' is outside the project or invalid."
    existed = os.path.isfile(p)
    backup = _checkpoint_copy(p) if existed else None
    try:
        os.makedirs(os.path.dirname(p), exist_ok=True)
        with open(p, "w", encoding="utf-8") as fh:
            fh.write(content)
    except OSError as exc:
        return f"Write failed: {exc}"
    _journal_add({"op": "update" if existed else "create", "path": rel, "backup": backup})
    return f"{'Updated' if existed else 'Created'} {rel} ({len(content)} chars)."


def _crud_mkdir(args: dict[str, Any]) -> str:
    guard = _writes_guard()
    if guard:
        return guard
    rel = str(args.get("path") or "").strip()
    p = _project_path(rel)
    if not p or p == os.path.realpath(PROJECT_ROOT):
        return f"Refused: invalid path '{rel}'."
    if os.path.exists(p):
        return f"'{rel}' already exists."
    try:
        os.makedirs(p)
    except OSError as exc:
        return f"Create-folder failed: {exc}"
    _journal_add({"op": "mkdir", "path": rel})
    return f"Created folder {rel}."


def _crud_move(args: dict[str, Any]) -> str:
    guard = _writes_guard()
    if guard:
        return guard
    src = str(args.get("from") or args.get("src") or "").strip()
    dst = str(args.get("to") or args.get("dest") or "").strip()
    ps, pd = _project_path(src), _project_path(dst)
    root = os.path.realpath(PROJECT_ROOT)
    if not ps or ps == root or not os.path.exists(ps):
        return f"Refused: source '{src}' is missing or invalid."
    if not pd or pd == root:
        return f"Refused: destination '{dst}' is invalid."
    if os.path.exists(pd):
        return f"Refused: '{dst}' already exists."
    try:
        os.makedirs(os.path.dirname(pd), exist_ok=True)
        os.rename(ps, pd)
    except OSError as exc:
        return f"Move failed: {exc}"
    _journal_add({"op": "move", "path": src, "to": dst})
    return f"Moved {src} → {dst}."


def _crud_delete(args: dict[str, Any]) -> str:
    guard = _writes_guard()
    if guard:
        return guard
    rel = str(args.get("path") or "").strip()
    p = _project_path(rel)
    root = os.path.realpath(PROJECT_ROOT)
    if not p or p == root or not os.path.exists(p):
        return f"Refused: '{rel}' is missing or invalid."
    trashname = uuid.uuid4().hex[:12] + "__" + os.path.basename(p.rstrip(os.sep))
    try:
        shutil.move(p, os.path.join(_agent_dir(), "trash", trashname))
    except OSError as exc:
        return f"Delete failed: {exc}"
    _journal_add({"op": "delete", "path": rel, "trash": trashname})
    return f"Deleted {rel} (moved to the project trash — recoverable via Undo)."


def _crud_undo(change_id: "str | None" = None) -> str:
    with _AGENT_LOCK:
        entries = _journal_load()
        target = None
        for e in reversed(entries):
            if e.get("undone"):
                continue
            if change_id is None or e.get("id") == change_id:
                target = e
                break
        if not target:
            return "Nothing to undo."
        op, rel, root = target.get("op"), target.get("path"), os.path.realpath(PROJECT_ROOT)
        try:
            if op == "create":
                p = _project_path(rel)
                if p and os.path.isfile(p):
                    os.remove(p)
            elif op == "update":
                p, bk = _project_path(rel), target.get("backup")
                src = os.path.join(_agent_dir(), "checkpoints", bk) if bk else None
                if p and src and os.path.isfile(src):
                    shutil.copy2(src, p)
            elif op == "mkdir":
                p = _project_path(rel)
                if p and os.path.isdir(p):
                    try:
                        os.rmdir(p)
                    except OSError:
                        pass
            elif op == "move":
                ps, pd = _project_path(rel), _project_path(target.get("to"))
                if ps and pd and os.path.exists(pd) and not os.path.exists(ps):
                    os.makedirs(os.path.dirname(ps), exist_ok=True)
                    os.rename(pd, ps)
            elif op == "delete":
                p, tn = _project_path(rel), target.get("trash")
                src = os.path.join(_agent_dir(), "trash", tn) if tn else None
                if p and src and os.path.exists(src) and not os.path.exists(p):
                    os.makedirs(os.path.dirname(p), exist_ok=True)
                    shutil.move(src, p)
        except OSError as exc:
            return f"Undo failed: {exc}"
        target["undone"] = True
        _journal_save(entries)
    return f"Undid {op} on {rel}."


# Tools the model can call (only listed when AGENT_WRITES_ENABLED). Registered in
# NATIVE_TOOLS so _dispatch_tool routes them for EVERY provider.
NATIVE_TOOLS["project_write"] = _crud_write
NATIVE_TOOLS["project_create_folder"] = _crud_mkdir
NATIVE_TOOLS["project_move"] = _crud_move
NATIVE_TOOLS["project_delete"] = _crud_delete

_CRUD_TOOL_SCHEMAS: list[dict[str, Any]] = [
    {"type": "function", "function": {
        "name": "project_write",
        "description": "Create a new file OR overwrite an existing one inside the project, "
                       "with the given full content. The change is journaled and undoable.",
        "parameters": {"type": "object", "properties": {
            "path": {"type": "string", "description": "project-relative path, e.g. 'src/app.py'"},
            "content": {"type": "string", "description": "the FULL new file content"}},
            "required": ["path", "content"]}}},
    {"type": "function", "function": {
        "name": "project_create_folder",
        "description": "Create a new folder (and any parents) inside the project. Journaled/undoable.",
        "parameters": {"type": "object", "properties": {
            "path": {"type": "string", "description": "project-relative folder path"}},
            "required": ["path"]}}},
    {"type": "function", "function": {
        "name": "project_move",
        "description": "Move or rename a file/folder within the project. Journaled/undoable.",
        "parameters": {"type": "object", "properties": {
            "from": {"type": "string", "description": "current project-relative path"},
            "to": {"type": "string", "description": "new project-relative path"}},
            "required": ["from", "to"]}}},
    {"type": "function", "function": {
        "name": "project_delete",
        "description": "Delete a file/folder in the project. It is moved to the project trash "
                       "(never hard-deleted), so it can be restored via Undo.",
        "parameters": {"type": "object", "properties": {
            "path": {"type": "string", "description": "project-relative path to delete"}},
            "required": ["path"]}}},
]


@app.get("/api/agent/writes")
def api_agent_writes_get() -> Response:
    return jsonify({"enabled": AGENT_WRITES_ENABLED})


@app.post("/api/agent/writes")
def api_agent_writes_set() -> Response:
    """Turn agent file edits on/off (off by default)."""
    global AGENT_WRITES_ENABLED
    data = request.get_json(silent=True) or {}
    AGENT_WRITES_ENABLED = bool(data.get("enabled"))
    return jsonify({"enabled": AGENT_WRITES_ENABLED})


@app.get("/api/agent/approval")
def api_agent_approval_get() -> Response:
    return jsonify({"required": AGENT_APPROVAL_REQUIRED})


@app.post("/api/agent/approval")
def api_agent_approval_set() -> Response:
    """Require Approve/Reject for each agent edit (on by default)."""
    global AGENT_APPROVAL_REQUIRED
    AGENT_APPROVAL_REQUIRED = bool((request.get_json(silent=True) or {}).get("required"))
    return jsonify({"required": AGENT_APPROVAL_REQUIRED})


@app.post("/api/agent/change/approve")
def api_change_approve() -> Response:
    """Apply a staged change (the only path that applies a model-proposed edit)."""
    cid = str((request.get_json(silent=True) or {}).get("id") or "")
    pend = _PENDING_CHANGES.pop(cid, None)
    if not pend:
        return jsonify({"error": "No such pending change (already handled or expired)."}), 404
    msg = _dispatch_tool(pend["name"], pend["args"])  # applies + journals via the CRUD core
    return jsonify({"ok": True, "message": msg})


@app.post("/api/agent/change/reject")
def api_change_reject() -> Response:
    cid = str((request.get_json(silent=True) or {}).get("id") or "")
    _PENDING_CHANGES.pop(cid, None)
    return jsonify({"ok": True})


@app.get("/api/fs/changes")
def api_fs_changes() -> Response:
    """The change journal (most recent first)."""
    return jsonify({"changes": list(reversed(_journal_load())), "enabled": AGENT_WRITES_ENABLED})


@app.post("/api/fs/undo")
def api_fs_undo() -> Response:
    """Undo a journaled change (the most recent, or a specific id)."""
    data = request.get_json(silent=True) or {}
    msg = _crud_undo(str(data["id"]) if data.get("id") else None)
    return jsonify({"message": msg})


@app.get("/api/fs/diff")
def api_fs_diff() -> Response:
    """A before→after unified diff for a journaled change, so the user can review
    exactly what the agent did. Uses the checkpoint copy (updates) or the trashed
    copy (deletes) saved at change time."""
    import difflib
    cid = request.args.get("id", "")
    entry = next((e for e in _journal_load() if e.get("id") == cid), None)
    if not entry:
        return jsonify({"error": "No such change."}), 404
    op, rel = entry.get("op"), entry.get("path")
    if op in ("mkdir", "move"):
        return jsonify({"op": op, "path": rel, "to": entry.get("to"), "diff": "",
                        "note": f"{op} — no file-content change to diff."})

    def _read(path: "str | None") -> str:
        if not path or not os.path.isfile(path):
            return ""
        try:
            with open(path, "rb") as fh:
                return fh.read(_FS_READ_CAP).decode("utf-8")
        except (OSError, UnicodeDecodeError):
            return ""

    cur = _project_path(rel)
    ck = os.path.join(_agent_dir(), "checkpoints", entry["backup"]) if entry.get("backup") else None
    tr = os.path.join(_agent_dir(), "trash", entry["trash"]) if entry.get("trash") else None
    before = _read(ck) if op == "update" else (_read(tr) if op == "delete" else "")
    after = _read(cur) if op in ("create", "update") else ""
    diff = "".join(difflib.unified_diff(
        before.splitlines(keepends=True), after.splitlines(keepends=True),
        fromfile=rel + " (before)", tofile=rel + " (after)", n=3))
    return jsonify({"op": op, "path": rel, "diff": diff})


# Direct (UI-driven) mutation endpoints — same core, gated the same way.
@app.post("/api/fs/write")
def api_fs_write() -> Response:
    return jsonify({"message": _crud_write(request.get_json(silent=True) or {})})


@app.post("/api/fs/mkdir")
def api_fs_mkdir() -> Response:
    return jsonify({"message": _crud_mkdir(request.get_json(silent=True) or {})})


@app.post("/api/fs/move")
def api_fs_move() -> Response:
    return jsonify({"message": _crud_move(request.get_json(silent=True) or {})})


@app.post("/api/fs/delete")
def api_fs_delete() -> Response:
    return jsonify({"message": _crud_delete(request.get_json(silent=True) or {})})


# ─────────────────────────────────────────────────────────────────────────
# Resumable sessions (Phase 6) — a server-side mirror of the browser's
# localStorage conversation store, so chats survive a cache clear and can be
# reopened. The browser stays the working store and POSTs each turn here; on
# startup it backfills any sessions missing locally. Stored as one JSON per
# session under .lls-sessions/ (gitignored).
# ─────────────────────────────────────────────────────────────────────────
SESSIONS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), ".lls-sessions")
os.makedirs(SESSIONS_DIR, exist_ok=True)
_SESSION_ID_RE = re.compile(r"^[A-Za-z0-9_-]{1,64}$")
_SESSION_MAX_BYTES = 8_000_000  # refuse absurdly large payloads


def _session_file(sid: Any) -> "str | None":
    if not isinstance(sid, str) or not _SESSION_ID_RE.match(sid):
        return None
    return os.path.join(SESSIONS_DIR, sid + ".json")


@app.get("/api/sessions")
def api_sessions_list() -> Response:
    """All stored sessions (full objects), newest first."""
    out: list[dict[str, Any]] = []
    try:
        for fn in os.listdir(SESSIONS_DIR):
            if not fn.endswith(".json"):
                continue
            try:
                with open(os.path.join(SESSIONS_DIR, fn), encoding="utf-8") as fh:
                    out.append(json.load(fh))
            except (OSError, ValueError):
                pass
    except OSError:
        pass
    out.sort(key=lambda s: s.get("updatedAt", 0), reverse=True)
    return jsonify({"sessions": out})


@app.post("/api/sessions")
def api_sessions_save() -> Response:
    """Create/replace one session (the browser mirrors each turn here)."""
    data = request.get_json(silent=True) or {}
    p = _session_file(str(data.get("id") or ""))
    if not p:
        return jsonify({"error": "Invalid or missing session id."}), 400
    blob = json.dumps(data)
    if len(blob) > _SESSION_MAX_BYTES:
        return jsonify({"error": "Session too large."}), 413
    try:
        with open(p, "w", encoding="utf-8") as fh:
            fh.write(blob)
    except OSError as exc:
        return jsonify({"error": str(exc)}), 500
    return jsonify({"ok": True, "id": data.get("id")})


@app.delete("/api/sessions/<sid>")
def api_sessions_delete(sid: str) -> Response:
    p = _session_file(sid)
    if not p:
        return jsonify({"error": "Invalid session id."}), 400
    try:
        if os.path.isfile(p):
            os.remove(p)
    except OSError as exc:
        return jsonify({"error": str(exc)}), 500
    return jsonify({"ok": True})


if __name__ == "__main__":
    import atexit
    import signal

    # Tear down the spawned MCP servers when this worker exits, so a Stop/restart
    # (or a crashed run) never leaks orphaned uv/python/node processes.
    atexit.register(MCP.close)

    def _graceful_exit(*_: Any) -> None:
        MCP.close()
        os._exit(0)

    for _sig in (signal.SIGTERM, signal.SIGINT):
        signal.signal(_sig, _graceful_exit)

    # Port comes from the supervisor (env PORT); defaults to 5050 when run
    # standalone (5000 collides with macOS AirPlay). use_reloader=False so the
    # supervisor's single subprocess maps to a single PID it can cleanly terminate.
    # threaded=True so streaming chats, racing Arena panes, and the (blocking) video
    # poll loop all run concurrently instead of serializing on one worker thread.
    port = int(os.environ.get("PORT", "5050"))
    app.run(host="127.0.0.1", port=port, debug=True, use_reloader=False, threaded=True)
